#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL_TYPE
import os

def create_keti_style_presentation():
    """KETI 스타일 25페이지 PPT 생성 (15분 발표용)"""
    
    # 기존 KETI PPT를 템플릿으로 사용
    template_path = '/home/kim/github-sdv/중국SDV표준 소개_KETI 박부식.pptx'
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        # 기존 슬라이드 제거 (템플릿 레이아웃만 사용)
        while len(prs.slides) > 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[0])
    else:
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
    
    # 색상 테마 정의 (KETI 스타일)
    KETI_BLUE = RGBColor(0, 82, 147)  # KETI 파란색
    KETI_NAVY = RGBColor(0, 32, 96)   # 진한 네이비
    KETI_LIGHT_BLUE = RGBColor(218, 238, 243)  # 연한 파란색
    KETI_GRAY = RGBColor(95, 96, 98)  # 회색
    
    # 슬라이드 1: 표지
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # KETI 로고 위치 (플레이스홀더)
    logo_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(14), Inches(0.5),
        Inches(1.5), Inches(0.8)
    )
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = KETI_BLUE
    logo_text = logo_box.text_frame
    logo_text.clear()
    p = logo_text.add_paragraph()
    p.text = "KETI"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "중국 SDV 표준 분석"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    p = title_frame.add_paragraph()
    p.text = "지능형 커넥티드카 서비스 인터페이스 표준"
    p.font.size = Pt(28)
    p.font.color.rgb = KETI_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 버전 정보
    version_box = slide.shapes.add_textbox(Inches(4), Inches(5), Inches(8), Inches(1))
    version_frame = version_box.text_frame
    version_frame.clear()
    
    p = version_frame.add_paragraph()
    p.text = "SDV/T 001-2022 Version 4 Beta 1"
    p.font.size = Pt(22)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 발표자 정보
    presenter_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(1))
    presenter_frame = presenter_box.text_frame
    presenter_frame.clear()
    
    p = presenter_frame.add_paragraph()
    p.text = "한국전자기술연구원 모빌리티플랫폼연구센터"
    p.font.size = Pt(16)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    p = presenter_frame.add_paragraph()
    p.text = "2025. 08. 26"
    p.font.size = Pt(14)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 2: 목차
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "Contents"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    
    # 목차 내용
    contents = [
        ("01", "SDV 패러다임 전환"),
        ("02", "중국 SDV 표준화 전략"),
        ("03", "표준 개발 참여사 현황"),
        ("04", "4-Layer 아키텍처 구조"),
        ("05", "Part 1: Atomic Service API"),
        ("06", "Part 2: Device Abstraction API"),
        ("07", "Version 4 주요 개선사항"),
        ("08", "글로벌 표준 비교 분석"),
        ("09", "한국 대응 전략 제안")
    ]
    
    y_pos = 2
    for num, title in contents:
        # 번호 박스
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(y_pos),
            Inches(0.8), Inches(0.6)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = KETI_BLUE
        
        text_frame = num_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 제목 텍스트
        text_box = slide.shapes.add_textbox(Inches(3), Inches(y_pos), Inches(10), Inches(0.6))
        text_frame = text_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = KETI_NAVY
        
        y_pos += 0.7
    
    # 슬라이드 3: SDV 개요
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    # 타이틀 밴드
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "01. SDV 패러다임 전환"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # SDV 정의
    def_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(7), Inches(3)
    )
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
    
    text_frame = def_box.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "Software Defined Vehicle"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\n• 소프트웨어가 차량 가치의 핵심\n• OTA를 통한 지속적 업데이트\n• 하드웨어와 소프트웨어 분리\n• 차량의 '스마트폰화'"
    p.font.size = Pt(18)
    p.line_spacing = 1.5
    
    # 시장 전망
    market_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(1.5),
        Inches(6.5), Inches(3)
    )
    market_box.fill.solid()
    market_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    
    text_frame = market_box.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "시장 전망"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\n• 2035년까지 CAGR 30-34%\n• 2030년 시장규모 6,500억 달러\n• 전체 차량가치의 60% 차지\n• 신규 수익모델 창출"
    p.font.size = Pt(18)
    p.line_spacing = 1.5
    
    # 차량 이미지 플레이스홀더
    car_image = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(5),
        Inches(8), Inches(3)
    )
    car_image.fill.solid()
    car_image.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    text_frame = car_image.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "[SDV 컨셉카 이미지]"
    p.font.size = Pt(16)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 4: 중국 전략
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    # 타이틀 밴드
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "02. 중국 SDV 표준화 전략"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # CAAM 소개
    caam_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(14), Inches(1.5)
    )
    caam_box.fill.solid()
    caam_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    
    text_frame = caam_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "중국자동차공업협회(CAAM) 소프트웨어 분과 주도"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "국가 차원의 통합 표준화 추진"
    p.font.size = Pt(18)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 전략 목표
    goals = [
        ("통합 표준", "산업 파편화 방지\n개발 효율성 극대화"),
        ("생태계 구축", "자체 기술 생태계\n글로벌 주도권 확보"),
        ("빠른 상용화", "개발 비용 절감\n시장 출시 가속화")
    ]
    
    x_pos = 1
    for title, desc in goals:
        goal_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(3.5),
            Inches(4.5), Inches(3)
        )
        goal_box.fill.solid()
        goal_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
        
        text_frame = goal_box.text_frame
        text_frame.clear()
        text_frame.margin_all = Inches(0.3)
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = KETI_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = "\n" + desc
        p.font.size = Pt(16)
        p.font.color.rgb = KETI_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.3
        
        x_pos += 5.2
    
    # 슬라이드 5: 참여사 (확장)
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    # 타이틀 밴드
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "03. 표준 개발 핵심 참여사"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # OEM 섹션
    oem_section = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(7.5), Inches(5)
    )
    oem_section.fill.solid()
    oem_section.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    text_frame = oem_section.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "완성차 업체 (OEM)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # OEM 로고 플레이스홀더들
    oems = ["BYD", "GWM", "Geely", "FAW", "SAIC", "Changan", "GAC", "Dongfeng"]
    for i, oem in enumerate(oems):
        x = Inches(1 + (i % 4) * 1.8)
        y = Inches(2.5 + (i // 4) * 1.5)
        
        logo_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(1.5), Inches(1)
        )
        logo_box.fill.solid()
        logo_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        text_frame = logo_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = f"[{oem}\nLogo]"
        p.font.size = Pt(12)
        p.font.color.rgb = KETI_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Tier-1 섹션
    tier1_section = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(1.5),
        Inches(7.5), Inches(5)
    )
    tier1_section.fill.solid()
    tier1_section.fill.fore_color.rgb = RGBColor(255, 248, 240)
    
    text_frame = tier1_section.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "부품사 및 IT 기업"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Tier-1 로고 플레이스홀더들
    tier1s = ["Huawei", "Bosch", "Continental", "Baidu", "Tencent", "Alibaba", "ZTE", "Xiaomi"]
    for i, tier1 in enumerate(tier1s):
        x = Inches(8.5 + (i % 4) * 1.8)
        y = Inches(2.5 + (i // 4) * 1.5)
        
        logo_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(1.5), Inches(1)
        )
        logo_box.fill.solid()
        logo_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        text_frame = logo_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = f"[{tier1}\nLogo]"
        p.font.size = Pt(12)
        p.font.color.rgb = KETI_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # 통계
    stat_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(6.8),
        Inches(8), Inches(1.5)
    )
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = KETI_BLUE
    
    text_frame = stat_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "60+ 기업 참여 | 520+ API 정의 | 6대 핵심 도메인"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 6: 4-Layer 아키텍처 (상세)
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    # 타이틀 밴드
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "04. SDV 4-Layer 아키텍처"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 계층 다이어그램
    layers = [
        ("애플리케이션 계층", "차량 앱, 사용자 서비스", RGBColor(100, 200, 100)),
        ("아토믹 서비스 계층", "Part 1: 기능 API (290+)", RGBColor(255, 200, 100)),
        ("디바이스 추상화 계층", "Part 2: 장치 API (230+)", RGBColor(200, 150, 255)),
        ("기초 플랫폼 계층", "OS, 하드웨어, 드라이버", RGBColor(150, 200, 255))
    ]
    
    for i, (name, desc, color) in enumerate(layers):
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(1.8 + i*1.4),
            Inches(8), Inches(1.2)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        
        text_frame = layer_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = f"{name}: {desc}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 화살표
        if i < 3:
            arrow_box = slide.shapes.add_textbox(Inches(5.8), Inches(3 + i*1.4), Inches(0.5), Inches(0.3))
            arrow_frame = arrow_box.text_frame
            arrow_frame.clear()
            p = arrow_frame.add_paragraph()
            p.text = "⬇"
            p.font.size = Pt(20)
            p.font.color.rgb = KETI_BLUE
            p.alignment = PP_ALIGN.CENTER
    
    # 특징 설명
    feature_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.5), Inches(1.8),
        Inches(4.5), Inches(5.4)
    )
    feature_box.fill.solid()
    feature_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
    
    text_frame = feature_box.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "핵심 특징"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    
    p = text_frame.add_paragraph()
    p.text = "\n✓ 계층 간 독립성\n\n✓ 표준화된 인터페이스\n\n✓ 하드웨어 추상화\n\n✓ 소프트웨어 재사용성\n\n✓ 개발 복잡도 감소"
    p.font.size = Pt(16)
    p.line_spacing = 1.5
    
    # 슬라이드 7-11: Part 1 Atomic Service API (5페이지로 확장)
    # 슬라이드 7: Part 1 개요
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "05. Part 1: Atomic Service API Overview"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 6대 도메인 헥사곤 배치
    domains = [
        ("BCM", "차체 제어", Inches(3), Inches(2)),
        ("TMS", "열 관리", Inches(7), Inches(2)),
        ("VCS", "운동 제어", Inches(11), Inches(2)),
        ("EMS", "에너지 관리", Inches(3), Inches(4.5)),
        ("ADAS", "첨단 운전자 보조", Inches(7), Inches(4.5)),
        ("HMI", "사용자 인터페이스", Inches(11), Inches(4.5))
    ]
    
    for code, name, x, y in domains:
        domain_box = slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON,
            x, y,
            Inches(3), Inches(2)
        )
        domain_box.fill.solid()
        domain_box.fill.fore_color.rgb = KETI_BLUE
        
        text_frame = domain_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = name
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # API 수 표시
    api_count_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(7),
        Inches(6), Inches(1)
    )
    api_count_box.fill.solid()
    api_count_box.fill.fore_color.rgb = RGBColor(255, 200, 100)
    
    text_frame = api_count_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "Total: 290+ APIs"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 8: BCM 도메인 상세
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "05-1. BCM Domain: Body Control Module"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # BCM 서비스 리스트
    services = [
        ("BCM_Door", "차문 제어", "unlock(), lock(), open(), close()"),
        ("BCM_Window", "창문 제어", "open(), close(), adjustPosition()"),
        ("BCM_Seat", "좌석 제어", "adjustPosition(), setHeating(), setVentilation()"),
        ("BCM_Light", "조명 제어", "turnOn(), turnOff(), setIntensity()"),
        ("BCM_Mirror", "미러 제어", "adjustAngle(), fold(), unfold()"),
        ("BCM_WiperWash", "와이퍼", "startWiping(), setSpeed(), spray()")
    ]
    
    y_pos = 1.8
    for service, desc, apis in services:
        # 서비스 박스
        service_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(y_pos),
            Inches(3), Inches(0.8)
        )
        service_box.fill.solid()
        service_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
        
        text_frame = service_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = service
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = KETI_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        # 설명
        desc_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(3), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.clear()
        p = desc_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = KETI_GRAY
        
        # API
        api_box = slide.shapes.add_textbox(Inches(8), Inches(y_pos), Inches(7), Inches(0.8))
        api_frame = api_box.text_frame
        api_frame.clear()
        p = api_frame.add_paragraph()
        p.text = apis
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = KETI_BLUE
        
        y_pos += 1
    
    # 슬라이드 9: TMS/VCS/EMS 도메인
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "05-2. TMS/VCS/EMS Domains"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 3개 도메인 박스
    domain_data = [
        ("TMS", "열 관리", ["AC 제어", "배터리 온도", "냉각수 관리", "히터 제어"], KETI_LIGHT_BLUE),
        ("VCS", "운동 제어", ["기어 변속", "브레이크", "조향", "서스펜션"], RGBColor(255, 240, 230)),
        ("EMS", "에너지 관리", ["충전 제어", "배터리 관리", "전력 분배", "에너지 최적화"], RGBColor(240, 255, 240))
    ]
    
    x_pos = 1
    for domain, title, items, color in domain_data:
        domain_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(2),
            Inches(4.8), Inches(5)
        )
        domain_box.fill.solid()
        domain_box.fill.fore_color.rgb = color
        
        text_frame = domain_box.text_frame
        text_frame.clear()
        text_frame.margin_all = Inches(0.3)
        
        p = text_frame.add_paragraph()
        p.text = domain
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = KETI_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(18)
        p.font.color.rgb = KETI_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = "\n"
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.line_spacing = 1.5
        
        x_pos += 5.2
    
    # 슬라이드 10: ADAS 도메인
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "05-3. ADAS Domain: Advanced Driver-Assistance"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # ADAS 센서 이미지 플레이스홀더
    sensor_image = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(7), Inches(3)
    )
    sensor_image.fill.solid()
    sensor_image.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    text_frame = sensor_image.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "[ADAS 센서 구성도]"
    p.font.size = Pt(16)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # ADAS 기능
    functions = [
        ("객체 인식", "차량, 보행자, 자전거"),
        ("차선 인식", "차선, 도로 경계"),
        ("신호 인식", "신호등, 표지판"),
        ("주차 보조", "주차 공간 탐지")
    ]
    
    x_pos = 8.5
    y_pos = 1.5
    for func, desc in functions:
        func_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(6.5), Inches(0.9)
        )
        func_box.fill.solid()
        func_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
        
        text_frame = func_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = f"{func}: {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = KETI_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1.1
    
    # 라이다 이미지
    lidar_image = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(5),
        Inches(7), Inches(2.5)
    )
    lidar_image.fill.solid()
    lidar_image.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    text_frame = lidar_image.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "[LiDAR Point Cloud 이미지]"
    p.font.size = Pt(14)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 11: 시나리오 예시
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "05-4. 시나리오: '퇴근 모드' 구현"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 시나리오 플로우
    steps = [
        ("1. 사용자 인증", "얼굴/지문 인식", "HMI_Auth.verify()"),
        ("2. 환경 설정", "시트/미러 조정", "BCM_Seat.adjust()"),
        ("3. 공조 제어", "온도/공기질 설정", "TMS_AC.setTemp()"),
        ("4. 엔터테인먼트", "음악/네비게이션", "HMI_Media.play()"),
        ("5. 주행 모드", "편안한 주행 설정", "VCS_Drive.setMode()")
    ]
    
    y_pos = 2
    for step, desc, api in steps:
        # 스텝 박스
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(1), Inches(y_pos),
            Inches(3.5), Inches(1)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = KETI_BLUE
        
        text_frame = step_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 설명
        desc_box = slide.shapes.add_textbox(Inches(5), Inches(y_pos), Inches(4), Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.clear()
        p = desc_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = KETI_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        # API
        api_box = slide.shapes.add_textbox(Inches(9.5), Inches(y_pos), Inches(5.5), Inches(1))
        api_frame = api_box.text_frame
        api_frame.clear()
        p = api_frame.add_paragraph()
        p.text = api
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = KETI_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1.2
    
    # 슬라이드 12-15: Part 2 Device Abstraction API (4페이지)
    # 슬라이드 12: Part 2 개요
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "06. Part 2: Device Abstraction API Overview"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 5대 도메인
    domains = [
        ("BCM", "차체 장치", Inches(2), Inches(3)),
        ("TMS", "열관리 장치", Inches(5.5), Inches(2)),
        ("PWT", "파워트레인", Inches(9), Inches(2)),
        ("CHS", "섀시", Inches(12.5), Inches(3)),
        ("ADAS", "센서", Inches(7), Inches(4.5))
    ]
    
    for code, name, x, y in domains:
        domain_box = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON,
            x, y,
            Inches(3), Inches(2)
        )
        domain_box.fill.solid()
        domain_box.fill.fore_color.rgb = RGBColor(200, 100, 100)
        
        text_frame = domain_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = name
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # API 수
    api_count_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(7),
        Inches(6), Inches(1)
    )
    api_count_box.fill.solid()
    api_count_box.fill.fore_color.rgb = RGBColor(200, 150, 255)
    
    text_frame = api_count_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "Total: 230+ APIs"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 13: 센서와 액추에이터
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "06-1. Sensors & Actuators"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 센서 섹션
    sensor_section = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(7), Inches(5.5)
    )
    sensor_section.fill.solid()
    sensor_section.fill.fore_color.rgb = KETI_LIGHT_BLUE
    
    text_frame = sensor_section.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "센서 (Sensors)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    sensors = [
        "• 온도 센서 (Temper)",
        "• 압력 센서 (Pressure)",
        "• 위치 센서 (Position)",
        "• 속도 센서 (Speed)",
        "• 카메라 (Camera)",
        "• 레이더 (Radar)",
        "• 라이다 (LiDAR)"
    ]
    
    for sensor in sensors:
        p = text_frame.add_paragraph()
        p.text = sensor
        p.font.size = Pt(14)
        p.line_spacing = 1.5
    
    # 액추에이터 섹션
    actuator_section = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(1.5),
        Inches(7), Inches(5.5)
    )
    actuator_section.fill.solid()
    actuator_section.fill.fore_color.rgb = RGBColor(255, 240, 230)
    
    text_frame = actuator_section.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "액추에이터 (Actuators)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    actuators = [
        "• 도어락 모터 (DoorLock)",
        "• 윈도우 모터 (WinMot)",
        "• 시트 모터 (SeatMot)",
        "• 펌프 (Pump)",
        "• 밸브 (Valve)",
        "• 릴레이 (Relay)",
        "• 컨버터 (Converter)"
    ]
    
    for actuator in actuators:
        p = text_frame.add_paragraph()
        p.text = actuator
        p.font.size = Pt(14)
        p.line_spacing = 1.5
    
    # 슬라이드 14: PWT 파워트레인 도메인
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "06-2. Powertrain Domain (전기차 중심)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 전기차 파워트레인 이미지
    ev_image = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(7), Inches(3)
    )
    ev_image.fill.solid()
    ev_image.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    text_frame = ev_image.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "[전기차 파워트레인 구성도]"
    p.font.size = Pt(16)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 주요 장치
    devices = [
        ("충전 포트", "Actr_ChrgElecLock", "충전 포트 잠금/해제"),
        ("배터리 관리", "Actr_HvBattCtrl", "고전압 배터리 제어"),
        ("DC-DC 컨버터", "Actr_DcdcCtrl", "전압 변환 제어"),
        ("온도 센서", "Snsr_BattTemp", "배터리 온도 모니터링")
    ]
    
    y_pos = 1.5
    for device, api, desc in devices:
        device_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.5), Inches(y_pos),
            Inches(6.5), Inches(1.2)
        )
        device_box.fill.solid()
        device_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
        
        text_frame = device_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = f"{device}\n{api}: {desc}"
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1.4
    
    # 배터리 이미지
    battery_image = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(5),
        Inches(7), Inches(2.5)
    )
    battery_image.fill.solid()
    battery_image.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    text_frame = battery_image.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "[고전압 배터리 팩 이미지]"
    p.font.size = Pt(14)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 15: API 연관 관계
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "06-3. API 계층 간 호출 흐름"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 호출 흐름도
    flow_data = [
        ("사용자 입력", "'창문 열기' 터치", KETI_LIGHT_BLUE),
        ("애플리케이션", "UI 이벤트 처리", RGBColor(240, 248, 255)),
        ("Atomic Service", "BCM_Window.open()", RGBColor(255, 240, 230)),
        ("Device Abstraction", "Actr_WinMot.setOper()", RGBColor(240, 255, 240)),
        ("하드웨어", "모터 회전", RGBColor(255, 255, 240))
    ]
    
    y_pos = 1.8
    for step, action, color in flow_data:
        flow_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3), Inches(y_pos),
            Inches(10), Inches(1)
        )
        flow_box.fill.solid()
        flow_box.fill.fore_color.rgb = color
        
        text_frame = flow_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = f"{step}: {action}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        if y_pos < 6:
            # 화살표
            arrow_box = slide.shapes.add_textbox(Inches(7.7), Inches(y_pos + 1), Inches(0.5), Inches(0.3))
            arrow_frame = arrow_box.text_frame
            arrow_frame.clear()
            p = arrow_frame.add_paragraph()
            p.text = "⬇"
            p.font.size = Pt(20)
            p.font.color.rgb = KETI_BLUE
            p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1.3
    
    # 슬라이드 16: Version 4 개선사항
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "07. Version 4 주요 개선사항"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 버전 비교
    version_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(14), Inches(1.5)
    )
    version_box.fill.solid()
    version_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    
    text_frame = version_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "Version 3 → Version 4 Beta 1"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # 개선사항
    improvements = [
        ("API 확장", "TMS 도메인: 4.9.23 → 4.9.29", "열관리 기능 강화"),
        ("성능 개선", "응답 시간 30% 단축", "실시간 처리 능력 향상"),
        ("호환성", "하드웨어 지원 확대", "더 많은 센서/액추에이터"),
        ("보안 강화", "인증/암호화 개선", "사이버보안 대응")
    ]
    
    y_pos = 3.5
    for title, change, desc in improvements:
        imp_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(y_pos),
            Inches(12), Inches(0.9)
        )
        imp_box.fill.solid()
        imp_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
        
        text_frame = imp_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = f"{title}: {change} - {desc}"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1.1
    
    # 슬라이드 17-18: 글로벌 비교 (2페이지)
    # 슬라이드 17: 표준 비교 테이블
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "08. 글로벌 SDV 표준 비교"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 비교 테이블 (시각적으로 개선)
    standards = [
        ("", "중국 SDV", "AUTOSAR", "SOAFEE"),
        ("주도", "CAAM", "컨소시엄", "Arm"),
        ("철학", "통합 표준", "안전 중심", "클라우드"),
        ("API 수", "520+", "400+", "200+"),
        ("강점", "빠른 상용화", "안정성", "확장성")
    ]
    
    # 헤더
    for i, header in enumerate(standards[0]):
        if i == 0:
            continue
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1 + i*4.5), Inches(1.5),
            Inches(4), Inches(0.8)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = KETI_BLUE
        
        text_frame = header_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = header
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # 데이터 행
    y_pos = 2.5
    for row in standards[1:]:
        # 행 제목
        row_title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(2), Inches(0.8)
        )
        row_title_box.fill.solid()
        row_title_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
        
        text_frame = row_title_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = row[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 데이터 셀
        for i, data in enumerate(row[1:]):
            data_box = slide.shapes.add_textbox(Inches(2.8 + i*4.5), Inches(y_pos), Inches(3.5), Inches(0.8))
            data_frame = data_box.text_frame
            data_frame.clear()
            p = data_frame.add_paragraph()
            p.text = data
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1
    
    # 슬라이드 18: 경쟁 환경 분석
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "08-1. SDV 표준 경쟁 구도"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 경쟁 구도 다이어그램
    # 중국
    china_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2),
        Inches(4.5), Inches(4)
    )
    china_box.fill.solid()
    china_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    
    text_frame = china_box.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "중국 SDV"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\n• 내수시장 활용\n• 국가 주도\n• 빠른 상용화\n• 생태계 구축"
    p.font.size = Pt(14)
    p.line_spacing = 1.5
    
    # 유럽
    europe_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.8), Inches(2),
        Inches(4.5), Inches(4)
    )
    europe_box.fill.solid()
    europe_box.fill.fore_color.rgb = RGBColor(240, 240, 255)
    
    text_frame = europe_box.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "유럽 (AUTOSAR)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 200)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\n• 안전성 중심\n• 표준 성숙도\n• OEM 협력\n• 글로벌 확산"
    p.font.size = Pt(14)
    p.line_spacing = 1.5
    
    # 미국
    us_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.5), Inches(2),
        Inches(4.5), Inches(4)
    )
    us_box.fill.solid()
    us_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    
    text_frame = us_box.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.3)
    
    p = text_frame.add_paragraph()
    p.text = "미국 (SOAFEE)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\n• 클라우드 네이티브\n• 오픈소스\n• IT 기업 주도\n• 혁신 중심"
    p.font.size = Pt(14)
    p.line_spacing = 1.5
    
    # 슬라이드 19-21: 한국 대응 전략 (3페이지)
    # 슬라이드 19: 현황 분석
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "09. 한국 SDV 현황 분석"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # SWOT 분석
    swot_data = [
        ("강점 (S)", ["완성차 기술력", "IT 인프라", "배터리 경쟁력"], KETI_LIGHT_BLUE),
        ("약점 (W)", ["생태계 부족", "표준화 지연", "규모의 경제"], RGBColor(255, 240, 230)),
        ("기회 (O)", ["신시장 진출", "기술 전환기", "정부 지원"], RGBColor(240, 255, 240)),
        ("위협 (T)", ["중국 표준", "기술 종속", "인재 부족"], RGBColor(255, 230, 230))
    ]
    
    x_pos = 1
    y_pos = 2
    for title, items, color in swot_data:
        swot_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(6.8), Inches(2.5)
        )
        swot_box.fill.solid()
        swot_box.fill.fore_color.rgb = color
        
        text_frame = swot_box.text_frame
        text_frame.clear()
        text_frame.margin_all = Inches(0.2)
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = KETI_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.line_spacing = 1.3
        
        if x_pos > 7:
            x_pos = 1
            y_pos += 2.8
        else:
            x_pos += 7.5
    
    # 슬라이드 20: K-SDV 전략 제안
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "09-1. K-SDV 표준 제안"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 전략 방향
    strategies = [
        ("Core + Extension", "핵심 API 표준화\n프로파일 확장", KETI_BLUE),
        ("국제 호환성", "AUTOSAR 연계\nISO 표준 준수", RGBColor(100, 200, 100)),
        ("보안 우선", "사이버보안 내장\nOTA 표준화", RGBColor(200, 100, 100)),
        ("융합 서비스", "MaaS 연계\nV2X 통합", RGBColor(200, 200, 100))
    ]
    
    x_pos = 0.5
    for title, desc, color in strategies:
        strat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(2),
            Inches(3.7), Inches(3)
        )
        strat_box.fill.solid()
        strat_box.fill.fore_color.rgb = color
        
        text_frame = strat_box.text_frame
        text_frame.clear()
        text_frame.margin_all = Inches(0.3)
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = f"\n{desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.3
        
        x_pos += 3.9
    
    # 추진 체계
    system_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(5.5),
        Inches(12), Inches(2)
    )
    system_box.fill.solid()
    system_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
    
    text_frame = system_box.text_frame
    text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = "추진 체계"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "정부 주도 + 산업계 협력 + 연구기관 지원 + 국제 협력"
    p.font.size = Pt(16)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 21: 로드맵
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "09-2. K-SDV 추진 로드맵"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 로드맵 타임라인
    timeline_data = [
        ("2025", "표준화 착수", "협의체 구성\nAPI 정의"),
        ("2026", "PoC 개발", "시범 적용\n검증"),
        ("2027", "상용화", "양산 적용\n생태계 확대"),
        ("2028+", "글로벌 확산", "해외 진출\n국제 표준화")
    ]
    
    x_pos = 1
    for year, title, desc in timeline_data:
        # 연도 원
        year_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x_pos), Inches(2),
            Inches(2), Inches(2)
        )
        year_circle.fill.solid()
        year_circle.fill.fore_color.rgb = KETI_BLUE
        
        text_frame = year_circle.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = year
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 설명 박스
        desc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos - 0.5), Inches(4.5),
            Inches(3), Inches(2)
        )
        desc_box.fill.solid()
        desc_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
        
        text_frame = desc_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = KETI_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        
        # 화살표
        if x_pos < 13:
            arrow = slide.shapes.add_textbox(Inches(x_pos + 2.2), Inches(2.8), Inches(1), Inches(0.5))
            arrow_frame = arrow.text_frame
            arrow_frame.clear()
            p = arrow_frame.add_paragraph()
            p.text = "→"
            p.font.size = Pt(24)
            p.font.color.rgb = KETI_GRAY
            p.alignment = PP_ALIGN.CENTER
        
        x_pos += 3.8
    
    # 슬라이드 22: 시사점
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "10. 종합 시사점"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 핵심 메시지
    key_msg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.8),
        Inches(14), Inches(2)
    )
    key_msg.fill.solid()
    key_msg.fill.fore_color.rgb = RGBColor(255, 250, 240)
    
    text_frame = key_msg.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "중국 SDV 표준은 기술 혁신과 생태계 주도권을 동시에 추구"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "글로벌 자동차 산업의 새로운 경쟁 축 형성"
    p.font.size = Pt(20)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # 시사점 리스트
    implications = [
        "✓ 표준 주도권이 미래 자동차 산업 주도권 결정",
        "✓ 소프트웨어 생태계 구축이 핵심 경쟁력",
        "✓ 한국은 독자 표준과 국제 협력의 균형 필요",
        "✓ 산학연관 협력 체계 구축 시급"
    ]
    
    y_pos = 4.5
    for imp in implications:
        imp_box = slide.shapes.add_textbox(Inches(3), Inches(y_pos), Inches(10), Inches(0.7))
        imp_frame = imp_box.text_frame
        imp_frame.clear()
        p = imp_frame.add_paragraph()
        p.text = imp
        p.font.size = Pt(18)
        p.font.color.rgb = KETI_NAVY
        
        y_pos += 0.8
    
    # 슬라이드 23: 향후 과제
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    title_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = KETI_NAVY
    title_band.line.fill.background()
    
    title_text = title_band.text_frame
    title_text.clear()
    title_text.margin_left = Inches(0.5)
    p = title_text.add_paragraph()
    p.text = "11. 향후 과제"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 과제 카테고리
    tasks = [
        ("기술 개발", ["API 상세 분석", "PoC 개발", "성능 검증"], KETI_LIGHT_BLUE),
        ("표준화", ["국내 표준 제정", "국제 표준 연계", "인증 체계"], RGBColor(255, 240, 230)),
        ("생태계", ["개발자 육성", "파트너십 구축", "오픈소스화"], RGBColor(240, 255, 240)),
        ("정책", ["규제 개선", "투자 확대", "인프라 구축"], RGBColor(255, 230, 230))
    ]
    
    x_pos = 0.5
    y_pos = 2
    for title, items, color in tasks:
        task_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(7.2), Inches(2.5)
        )
        task_box.fill.solid()
        task_box.fill.fore_color.rgb = color
        
        text_frame = task_box.text_frame
        text_frame.clear()
        text_frame.margin_all = Inches(0.2)
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = KETI_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.line_spacing = 1.2
        
        if x_pos > 7:
            x_pos = 0.5
            y_pos += 2.8
        else:
            x_pos += 7.8
    
    # 슬라이드 24: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    # 배경 그라데이션 효과
    bg_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(9)
    )
    bg_box.fill.solid()
    bg_box.fill.fore_color.rgb = KETI_NAVY
    bg_box.line.fill.background()
    
    # Q&A 텍스트
    qa_box = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(10), Inches(2))
    qa_frame = qa_box.text_frame
    qa_frame.clear()
    
    p = qa_frame.add_paragraph()
    p.text = "Q & A"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 감사 메시지
    thanks_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(10), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.clear()
    
    p = thanks_frame.add_paragraph()
    p.text = "감사합니다"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 25: 연락처
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    
    # KETI 로고 플레이스홀더
    logo_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.5), Inches(1),
        Inches(3), Inches(1.5)
    )
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = KETI_BLUE
    
    text_frame = logo_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "KETI"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 연락처 정보
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(3),
        Inches(10), Inches(4)
    )
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = KETI_LIGHT_BLUE
    
    text_frame = contact_box.text_frame
    text_frame.clear()
    text_frame.margin_all = Inches(0.5)
    
    p = text_frame.add_paragraph()
    p.text = "한국전자기술연구원"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = KETI_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\n모빌리티플랫폼연구센터"
    p.font.size = Pt(22)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\n경기도 성남시 분당구 대왕판교로 712번길 22"
    p.font.size = Pt(16)
    p.font.color.rgb = KETI_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "\nwww.keti.re.kr"
    p.font.size = Pt(18)
    p.font.color.rgb = KETI_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return prs

# 프레젠테이션 생성 및 저장
prs = create_keti_style_presentation()
prs.save('/home/kim/github-sdv/China_SDV_Standard_KETI_Style_25pages.pptx')
print("KETI 스타일 25페이지 프레젠테이션이 생성되었습니다: China_SDV_Standard_KETI_Style_25pages.pptx")
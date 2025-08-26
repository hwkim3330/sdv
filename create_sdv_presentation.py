#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    
    title_shape.text = title
    if subtitle_shape:
        subtitle_shape.text = subtitle
    
    return slide

def add_content_slide(prs, title, content):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if content_shape:
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            
    return slide

def add_two_content_slide(prs, title, left_title, left_content, right_title, right_content):
    """두 컬럼 슬라이드 추가"""
    slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # 왼쪽 컬럼
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(4.5)
    
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.clear()
    
    p = left_frame.add_paragraph()
    p.text = left_title
    p.font.bold = True
    p.font.size = Pt(20)
    
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        
    # 오른쪽 컬럼
    right_left = Inches(5.5)
    right_box = slide.shapes.add_textbox(right_left, top, width, height)
    right_frame = right_box.text_frame
    right_frame.clear()
    
    p = right_frame.add_paragraph()
    p.text = right_title
    p.font.bold = True
    p.font.size = Pt(20)
    
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
    
    return slide

def create_presentation():
    """중국 SDV 표준 분석 프레젠테이션 생성"""
    prs = Presentation()
    
    # 슬라이드 1: 표지
    add_title_slide(prs, 
                    "중국 지능형 커넥티드카 서비스 인터페이스 표준\n(SDV/T 001-2022) 분석 보고",
                    "Part 1. Atomic Service & Part 2. Device Abstraction API 중심\nVersion 4 Beta 1")
    
    # 슬라이드 2: SDV 전환과 표준의 중요성
    add_content_slide(prs,
                      "SDV 전환과 표준의 중요성",
                      ["• 소프트웨어 중심 자동차(SDV): 하드웨어가 아닌 소프트웨어에 의해 기능이 정의",
                       "• 지속적 업데이트: 스마트폰처럼 SW 업데이트만으로 새로운 기능 추가",
                       "• 표준 API의 역할: 애플리케이션-소프트웨어-하드웨어 간 표준화된 인터페이스",
                       "• 개발 복잡성 감소 및 산업 생태계 구축의 기반"])
    
    # 슬라이드 3: 중국 SDV 표준화 동향
    add_content_slide(prs,
                      "중국의 표준화 전략: SDV/T 001-2022",
                      ["• 주도 기관: 중국자동차공업협회(CAAM) 소프트웨어 분과",
                       "• 표준화 목표:",
                       "  - 중국 내 산업 파편화 방지 및 통일된 개발 규격 제공",
                       "  - 자체 기술 생태계 구축으로 글로벌 SDV 시장 주도권 확보",
                       "  - 부품사-완성차 간 개발 비용 절감 및 빠른 상용화"])
    
    # 슬라이드 4: 표준 개발 핵심 참여사
    add_two_content_slide(prs,
                          "표준 개발 핵심 참여사",
                          "주요 완성차 업체(OEM)",
                          ["BYD(比亚迪汽车工业有限公司)",
                           "GWM(长城汽车股份有限公司)",
                           "Geely(浙江吉利汽车研究总院)",
                           "FAW(中国第一汽车集团有限公司)",
                           "SAIC(上海汽车集团股份有限公司)"],
                          "주요 부품사 및 SW 기업",
                          ["Huawei(华为技术有限公司)",
                           "Bosch(博世汽车部件)",
                           "Continental(大陆投资)",
                           "Baidu(北京百度智行科技有限公司)"])
    
    # 슬라이드 5: 4-Layer 아키텍처
    add_content_slide(prs,
                      "SDV 서비스 소프트웨어 4계층 아키텍처",
                      ["• 애플리케이션 계층: 사용자 경험 및 차량 특화 기능 구현",
                       "• 아토믹 서비스 계층 (Part 1): 표준화된 기능 단위 제공",
                       "• 디바이스 추상화 계층 (Part 2): 하드웨어 제어 인터페이스 표준화",
                       "• 기초 플랫폼 계층: OS, 컴퓨팅 하드웨어 등 기본 실행 환경"])
    
    # 슬라이드 6: API 계층의 역할과 관계
    add_content_slide(prs,
                      "소프트웨어와 하드웨어의 분리(Decoupling)",
                      ["• 아토믹 서비스 API (Part 1):",
                       "  - 애플리케이션 개발자가 물리적 구조를 몰라도 기능 호출 가능",
                       "  - 예: '창문을 열어줘' 같은 기능을 쉽게 호출",
                       "",
                       "• 디바이스 추상화 API (Part 2):",
                       "  - 특정 제조사 부품에 종속되지 않은 표준화된 장치 제어",
                       "  - 예: '모터를 정방향으로 회전' 같은 표준 명령 사용",
                       "",
                       "• 기대 효과: 하드웨어 교체 시 SW 변경 최소화, 재사용성 극대화"])
    
    # 슬라이드 7: Part 1 개요
    add_content_slide(prs,
                      "Part 1: 아토믹 서비스 API - '기능'의 표준화",
                      ["• 애플리케이션이 차량 핵심 기능을 쉽게 사용하도록 표준화된 서비스 집합",
                       "",
                       "• 6개 핵심 도메인:",
                       "  - BCM (Body Control Module): 차체 제어",
                       "  - TMS (Thermal Management System): 열 관리",
                       "  - VCS (Vehicle Control System): 차량 운동 제어",
                       "  - EMS (Energy Management System): 에너지 관리",
                       "  - ADAS (Advanced Driver-Assistance Systems): 첨단 운전자 보조",
                       "  - HMI (Human Machine Interface): 사용자 인터페이스"])
    
    # 슬라이드 8: BCM 도메인
    add_content_slide(prs,
                      "BCM: 차량 편의 기능 서비스",
                      ["• BCM_Door: unlock(), lock(), open(), close(), adjustPosition()",
                       "• BCM_Window: lock(), unlock(), open(), close(), adjustPosition()",
                       "• BCM_Seat: adjustMainXDir(), adjustBackRestAngle()",
                       "• BCM_Light: turnOn(), turnOff() (브레이크등, 방향지시등, 전조등)",
                       "• BCM_WiperWash: startWiping(), stopWiping(), startSprayWashing()"])
    
    # 슬라이드 9: TMS/VCS/EMS 도메인
    add_content_slide(prs,
                      "TMS/VCS/EMS: 차량 운행 핵심 기능 서비스",
                      ["• TMS (열 관리):",
                       "  - TMS_AC.setTargetTemp(): 실내 목표 온도 설정",
                       "  - TMS_Battery.setTargetTemp(): 배터리 목표 온도 설정",
                       "",
                       "• VCS (차량 운동 제어):",
                       "  - VCS_Gear.setTarget(): 목표 기어 설정",
                       "  - VCS_Brake.setTargetAx(): 목표 감속도 설정",
                       "",
                       "• EMS (에너지 관리):",
                       "  - EMS_Charging.start(): 충전 시작",
                       "  - EMS_HVBatt.getSOC(): 고전압 배터리 충전 상태 조회"])
    
    # 슬라이드 10: ADAS 도메인
    add_content_slide(prs,
                      "ADAS: 지능형 주행 기능 서비스",
                      ["• ADAS_Perception (시각 인지):",
                       "  - getTrackObjects(): 차량, 보행자 등 객체 정보 획득",
                       "  - getLaneline(): 차선 정보 획득",
                       "  - getTrafficLight(): 신호등 정보 획득",
                       "  - getParkingSpace(): 주차 공간 정보 획득",
                       "",
                       "• ADAS_Radar / Lidar (레이더/라이다 인지):",
                       "  - notifyObjects(): 레이더/라이다가 탐지한 객체 정보 획득"])
    
    # 슬라이드 11: 시나리오 예시
    add_content_slide(prs,
                      "시나리오: '퇴근 모드' 기능 구현",
                      ["• 시트 및 미러 조정:",
                       "  - BCM_Seat.adjustMainXDir(position): 운전석 시트 위치 조정",
                       "  - BCM_RearView.adjustXAngle(angle): 사이드미러 각도 조정",
                       "",
                       "• 실내 환경 설정:",
                       "  - TMS_AC.setTargetTemp(22.0): 실내 온도 22도로 설정",
                       "  - TMS_Purifier.turnOn(): 공기 청정 기능 활성화",
                       "",
                       "• 안마 기능 활성화:",
                       "  - BCM_Massage.startMassage(KNEAD, GENTLE): 부드러운 주무르기"])
    
    # 슬라이드 12: Part 2 개요
    add_content_slide(prs,
                      "Part 2: 디바이스 추상화 API - '장치'의 표준화",
                      ["• 차량의 물리적 장치(센서, 모터 등) 제어를 위한 표준 인터페이스",
                       "• 특정 제조사 부품에 상관없이 동일한 API로 제어",
                       "",
                       "• 5개 핵심 도메인:",
                       "  - BCM: 차체",
                       "  - TMS: 열 관리",
                       "  - PWT (Powertrain): 파워트레인",
                       "  - CHS (Chassis): 섀시",
                       "  - ADAS: 첨단 운전자 보조"])
    
    # 슬라이드 13: BCM 장치 제어
    add_content_slide(prs,
                      "BCM: 차체 편의 장치 제어",
                      ["• Actr_DoorLock: 도어 잠금 모터 제어",
                       "• Actr_DoubleHallMot: 홀 센서가 장착된 창문/선루프 모터 제어",
                       "• Snsr_WinSwt: 윈도우 스위치 신호 입력",
                       "• Actr_SeatHeatr: 시트 열선 제어",
                       "• Snsr_SeatOccupied: 좌석 탑승 감지 센서"])
    
    # 슬라이드 14: PWT 도메인
    add_content_slide(prs,
                      "PWT: 동력 및 충전 관련 장치 제어",
                      ["• Actr_ChrgElecLock: 충전 포트 잠금 장치 제어",
                       "• Snsr_AcChrgPortT: 교류(AC) 충전 포트 온도 센서",
                       "• Actr_HvBattCtrl: 고전압 배터리 릴레이 제어",
                       "• Snsr_HvBattCellInfo: 고전압 배터리 셀 전압/온도 조회",
                       "• Actr_DcdcCtrl: DC-DC 컨버터 제어"])
    
    # 슬라이드 15: ADAS 센서
    add_content_slide(prs,
                      "ADAS: 센서 원시 데이터(Raw Data) 제공",
                      ["• Snsr_Camera:",
                       "  - ntfRawData: 카메라 원본 이미지",
                       "  - ntfEncodeData: 인코딩된 비디오 스트림",
                       "",
                       "• Snsr_Radar: 밀리미터파 레이더 반사파 정보",
                       "• Snsr_Lidar: 라이다 포인트 클라우드 데이터",
                       "• Snsr_IMU: 관성 측정 장치(가속도, 각속도) 데이터",
                       "• Snsr_GPS: GPS 위치 정보"])
    
    # 슬라이드 16: API 연관 관계
    add_content_slide(prs,
                      "API 계층 간 호출 흐름: '창문 열기'",
                      ["1. 애플리케이션: 사용자가 '창문 열기' 버튼 터치",
                       "",
                       "2. 아토믹 서비스 API:",
                       "   - BCM_Window.open() API 호출",
                       "",
                       "3. 서비스 내부 로직:",
                       "   - 창문을 열기 위해 필요한 물리적 장치 제어",
                       "",
                       "4. 디바이스 추상화 API:",
                       "   - Actr_DoubleHallMot.setOper(dir=UPWARD, dutyRat=100)",
                       "   - 윈도우 모터를 정방향으로 100% 출력 구동",
                       "",
                       "5. 하드웨어 피드백:",
                       "   - ntfHallQuarterCnt()로 현재 위치 전달"])
    
    # 슬라이드 17: 특징 요약
    add_content_slide(prs,
                      "중국 SDV 표준의 특징 요약",
                      ["• 생태계:",
                       "  - BYD, Huawei 등 중국 완성차, 부품사, SW 기업 대거 참여",
                       "  - 강력한 산업 생태계 기반 표준 제정",
                       "",
                       "• 구조:",
                       "  - '기능(Atomic Service)'과 '장치(Device Abstraction)' 명확히 분리",
                       "  - 소프트웨어 유연성과 재사용성 극대화",
                       "",
                       "• 포괄성:",
                       "  - 편의 기능부터 ADAS, 파워트레인까지 차량 전체 시스템 포괄"])
    
    # 슬라이드 18: 시사점 및 제언
    add_content_slide(prs,
                      "시사점 및 국내 전략 제언",
                      ["• 시사점:",
                       "  - 중국 시장 진출 시 표준 준수가 필수 요건화 가능성",
                       "  - 중국 중심 SDV 기술 생태계 빠른 확장으로 경쟁 심화",
                       "",
                       "• 국내 전략 제언:",
                       "  - K-SDV 표준화 협의체 가속화",
                       "  - 국내 완성차 및 부품사 중심 기술 생태계 보호/발전",
                       "  - 표준 API 기반 서비스/애플리케이션 개발 역량 강화",
                       "  - 고부가가치 소프트웨어 시장 대비"])
    
    # 슬라이드 19: 한국형 표준 제안
    add_content_slide(prs,
                      "한국형 SDV 표준 제안 방향",
                      ["• Core API + Extension Profiles:",
                       "  - 핵심 API만 표준화, 확장은 Profile 형식으로 추가",
                       "",
                       "• 국제 표준 호환성:",
                       "  - VISS(W3C), ISO 20078, AUTOSAR과 매핑 테이블 제공",
                       "",
                       "• 보안/OTA 기본 내장:",
                       "  - API 호출 시 인증/권한/암호화 필수",
                       "  - OTA/SW 모듈 업그레이드 관련 API 포함",
                       "",
                       "• 차량-서비스 융합 API:",
                       "  - MaaS, C-ITS, 클라우드 연계 API 정의",
                       "  - V2X 및 데이터 공유 기능 표준화"])
    
    # 슬라이드 20: 향후 계획
    add_content_slide(prs,
                      "향후 추진 계획",
                      ["• 도메인별 심층 분석:",
                       "  - 핵심 도메인(ADAS, PWT 등) API 상세 사양 분석",
                       "  - PoC(개념증명) 추진",
                       "",
                       "• 국내 표준화 연계:",
                       "  - 국내 표준화 협의체에 중국 표준 분석 결과 공유",
                       "  - 주요 OEM 및 부품사와 협력",
                       "  - 국내 실정에 맞는 표준 API 개발 추진"])
    
    # 슬라이드 21: Q&A
    add_title_slide(prs, "Q & A", "감사합니다.")
    
    return prs

# 프레젠테이션 생성 및 저장
prs = create_presentation()
prs.save('/home/kim/github-sdv/중국_SDV_표준_분석_v4.pptx')
print("프레젠테이션이 성공적으로 생성되었습니다: 중국_SDV_표준_분석_v4.pptx")
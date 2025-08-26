#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
중국 SDV 표준 집중 프레젠테이션 - SDV/T 001-2022 Version 4 Beta 1
최신 동향과 API 상세 분석 포함
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 프레젠테이션 생성 (16:9 비율)
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 색상 정의 (중국 국가 색상과 KETI 색상 조합)
CHINA_RED = RGBColor(238, 28, 37)
CHINA_GOLD = RGBColor(255, 215, 0)
DARK_BLUE = RGBColor(0, 32, 96)
LIGHT_BLUE = RGBColor(218, 238, 243)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK_GRAY = RGBColor(64, 64, 64)

def add_title_slide():
    """타이틀 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경 그라데이션 효과를 위한 상단 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = CHINA_RED
    top_bar.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(2))
    title = title_box.text_frame
    title.text = "중국 SDV 표준 상세 분석"
    title.paragraphs[0].font.size = Pt(56)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    title.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 서브타이틀
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(1.5))
    subtitle = subtitle_box.text_frame
    subtitle.text = "SDV/T 001-2022 Version 4 Beta 1\n中国智能网联汽车服务接口规范"
    for para in subtitle.paragraphs:
        para.font.size = Pt(28)
        para.font.color.rgb = DARK_BLUE
        para.alignment = PP_ALIGN.CENTER
    
    # 버전 정보
    version_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
    version = version_box.text_frame
    version.text = "Part 1: 원자 서비스 API (290+ APIs)\nPart 2: 디바이스 추상화 API (230+ APIs)"
    for para in version.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = DARK_GRAY
        para.alignment = PP_ALIGN.CENTER
    
    # 발표 정보
    info_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(14), Inches(1.5))
    info = info_box.text_frame
    info.text = "한국전자기술연구원 (KETI)\n2025년 1월"
    for para in info.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = GRAY
        para.alignment = PP_ALIGN.CENTER

def add_content_slide(title_text, content_items, slide_number, code_example=None):
    """컨텐츠 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 상단 타이틀 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(14), Inches(0.8))
    title = title_box.text_frame
    title.text = title_text
    title.paragraphs[0].font.size = Pt(32)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    
    # 컨텐츠 영역
    content_top = Inches(1.5)
    
    if code_example:
        # 코드 예시가 있는 경우
        for item in content_items:
            item_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(7), Inches(0.6))
            item_text = item_box.text_frame
            item_text.text = f"• {item}"
            item_text.paragraphs[0].font.size = Pt(18)
            item_text.paragraphs[0].font.color.rgb = DARK_GRAY
            content_top += Inches(0.6)
        
        # 코드 박스
        code_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8.5), Inches(1.8),
            Inches(7), Inches(5.5)
        )
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        # 코드 텍스트
        code_text_box = slide.shapes.add_textbox(Inches(8.7), Inches(2), Inches(6.6), Inches(5))
        code_text = code_text_box.text_frame
        code_text.text = code_example
        code_text.paragraphs[0].font.name = "Consolas"
        code_text.paragraphs[0].font.size = Pt(12)
        code_text.paragraphs[0].font.color.rgb = BLACK
    else:
        # 일반 컨텐츠
        for item in content_items:
            if isinstance(item, dict):
                if 'heading' in item:
                    heading_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(14), Inches(0.6))
                    heading = heading_box.text_frame
                    heading.text = item['heading']
                    heading.paragraphs[0].font.size = Pt(24)
                    heading.paragraphs[0].font.bold = True
                    heading.paragraphs[0].font.color.rgb = CHINA_RED
                    content_top += Inches(0.7)
                
                if 'bullets' in item:
                    for bullet in item['bullets']:
                        bullet_box = slide.shapes.add_textbox(Inches(1.5), content_top, Inches(13), Inches(0.5))
                        bullet_text = bullet_box.text_frame
                        bullet_text.text = f"• {bullet}"
                        bullet_text.paragraphs[0].font.size = Pt(18)
                        bullet_text.paragraphs[0].font.color.rgb = DARK_GRAY
                        content_top += Inches(0.55)
            else:
                item_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(14), Inches(0.6))
                item_text = item_box.text_frame
                item_text.text = f"• {item}"
                item_text.paragraphs[0].font.size = Pt(20)
                item_text.paragraphs[0].font.color.rgb = DARK_GRAY
                content_top += Inches(0.65)
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    page = page_box.text_frame
    page.text = str(slide_number)
    page.paragraphs[0].font.size = Pt(12)
    page.paragraphs[0].font.color.rgb = GRAY
    page.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_api_table_slide(title_text, table_data, slide_number):
    """API 테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(14), Inches(0.8))
    title = title_box.text_frame
    title.text = title_text
    title.paragraphs[0].font.size = Pt(32)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    
    # 테이블
    rows = len(table_data)
    cols = len(table_data[0])
    
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(15), Inches(6)).table
    
    # 테이블 스타일링
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(table_data[i][j])
            
            if i == 0:  # 헤더 행
                cell.fill.solid()
                cell.fill.fore_color.rgb = CHINA_RED
                para = cell.text_frame.paragraphs[0]
                para.font.color.rgb = WHITE
                para.font.bold = True
                para.font.size = Pt(16)
            else:
                para = cell.text_frame.paragraphs[0]
                para.font.color.rgb = BLACK
                para.font.size = Pt(14)
            
            para.alignment = PP_ALIGN.CENTER
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    page = page_box.text_frame
    page.text = str(slide_number)
    page.paragraphs[0].font.size = Pt(12)
    page.paragraphs[0].font.color.rgb = GRAY

def add_architecture_slide(slide_number):
    """아키텍처 다이어그램 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(14), Inches(0.8))
    title = title_box.text_frame
    title.text = "중국 SDV 표준 4계층 아키텍처"
    title.paragraphs[0].font.size = Pt(32)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    
    # Layer 1 - 애플리케이션
    layer1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(1.8),
        Inches(12), Inches(1.2)
    )
    layer1.fill.solid()
    layer1.fill.fore_color.rgb = RGBColor(255, 230, 230)
    layer1_text = layer1.text_frame
    layer1_text.text = "Layer 1: 애플리케이션 계층\nOEM Apps | 3rd Party Apps | User Services"
    for para in layer1_text.paragraphs:
        para.font.size = Pt(16)
        para.alignment = PP_ALIGN.CENTER
    
    # Layer 2 - 아토믹 서비스
    layer2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(3.2),
        Inches(12), Inches(1.2)
    )
    layer2.fill.solid()
    layer2.fill.fore_color.rgb = RGBColor(230, 255, 230)
    layer2_text = layer2.text_frame
    layer2_text.text = "Layer 2: 아토믹 서비스 API (Part 1)\nBCM | TMS | VCS | EMS | ADAS | HMI (290+ APIs)"
    for para in layer2_text.paragraphs:
        para.font.size = Pt(16)
        para.alignment = PP_ALIGN.CENTER
    
    # Layer 3 - 디바이스 추상화
    layer3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(4.6),
        Inches(12), Inches(1.2)
    )
    layer3.fill.solid()
    layer3.fill.fore_color.rgb = RGBColor(230, 230, 255)
    layer3_text = layer3.text_frame
    layer3_text.text = "Layer 3: 디바이스 추상화 API (Part 2)\nActuators | Sensors | ECUs (230+ APIs)"
    for para in layer3_text.paragraphs:
        para.font.size = Pt(16)
        para.alignment = PP_ALIGN.CENTER
    
    # Layer 4 - 기초 플랫폼
    layer4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(6),
        Inches(12), Inches(1.2)
    )
    layer4.fill.solid()
    layer4.fill.fore_color.rgb = RGBColor(240, 240, 240)
    layer4_text = layer4.text_frame
    layer4_text.text = "Layer 4: 기초 플랫폼 계층\nLinux | QNX | Android Automotive | RTOS"
    for para in layer4_text.paragraphs:
        para.font.size = Pt(16)
        para.alignment = PP_ALIGN.CENTER
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    page = page_box.text_frame
    page.text = str(slide_number)
    page.paragraphs[0].font.size = Pt(12)
    page.paragraphs[0].font.color.rgb = GRAY

# 슬라이드 생성
slide_num = 1

# 1. 타이틀 슬라이드
add_title_slide()
slide_num += 1

# 2. 목차
add_content_slide(
    "목차",
    [
        "1. 중국 SDV 표준화 현황 (2024-2025)",
        "2. CAAM 소프트웨어 분과 활동",
        "3. SDV/T 001-2022 표준 개요",
        "4. Part 1: 아토믹 서비스 API 상세",
        "5. Part 2: 디바이스 추상화 API 상세",
        "6. Version 3 → Version 4 변경사항",
        "7. API 구현 예시 및 코드",
        "8. 중국 주요 기업 구현 현황",
        "9. 글로벌 영향력 및 시장 전망",
        "10. 한국의 대응 방향"
    ],
    slide_num
)
slide_num += 1

# 3. 중국 SDV 표준화 최신 현황
add_content_slide(
    "중국 SDV 표준화 최신 현황 (2024-2025)",
    [
        {
            'heading': '2024년 주요 이정표',
            'bullets': [
                'CAAM 소프트웨어 분과: 393개 아토믹 서비스 API 공개',
                '269개 디바이스 추상화 API 업계 공개',
                '2024년 11월 중국 자동차 소프트웨어 컨퍼런스 개최',
                'Version 4 Beta 1 산업계 적용 확대'
            ]
        },
        {
            'heading': '2025년 전망',
            'bullets': [
                '중국 NEV 판매 1,600만대 예상 (YoY 24.4%↑)',
                'SDV 시장 2030년까지 1.23조 달러 규모 성장',
                'L3 자율주행 차량 30% 달성 목표',
                '차량 컴퓨팅 파워 5,000 TOPS 초과'
            ]
        }
    ],
    slide_num
)
slide_num += 1

# 4. CAAM 소프트웨어 분과
add_content_slide(
    "CAAM 소프트웨어 분과 주도 표준화",
    [
        {
            'heading': '조직 구성',
            'bullets': [
                '위원장: TAN Minqiang (중국자동차혁신공사 CEO)',
                '참여 기업: 60개 이상 (OEM, Tier-1, IT 기업)',
                '2022년 3월 30일: API Reference Specification 2.0 발표',
                '반년 주기 업데이트 (Beta → 정식 버전)'
            ]
        },
        {
            'heading': '표준화 전략',
            'bullets': [
                '오픈소스 기반 생태계 구축',
                '하드웨어-소프트웨어 완전 분리',
                '도메인별 API 표준화',
                '국제 표준과의 선택적 호환성'
            ]
        }
    ],
    slide_num
)
slide_num += 1

# 5. SDV/T 001-2022 표준 구조
add_architecture_slide(slide_num)
slide_num += 1

# 6. Part 1: 아토믹 서비스 API 개요
add_api_table_slide(
    "Part 1: 아토믹 서비스 API - 6대 도메인 (290+ APIs)",
    [
        ["도메인", "서비스 수", "대표 API", "주요 기능"],
        ["BCM\n(차체제어)", "31개", "BCM_Door, BCM_Window\nBCM_Seat, BCM_Light", "도어/창문/시트/조명 제어"],
        ["TMS\n(열관리)", "8개", "TMS_AC, TMS_Battery\nTMS_Device", "공조/배터리 열관리"],
        ["VCS\n(차량제어)", "12개", "VCS_Gear, VCS_Brake\nVCS_Steering", "기어/브레이크/조향"],
        ["EMS\n(에너지)", "15개", "EMS_Charging, EMS_HVBatt\nEMS_PowerDist", "충전/배터리/전력분배"],
        ["ADAS\n(지능형)", "18개", "ADAS_Perception\nADAS_Fusion", "인지/센서융합"],
        ["HMI\n(인터페이스)", "10개", "HMI_Display, HMI_Audio\nHMI_Navigation", "디스플레이/오디오/내비"]
    ],
    slide_num
)
slide_num += 1

# 7. BCM 도메인 API 상세
add_content_slide(
    "BCM 도메인 API 구현 예시",
    [
        "BCM_Door: 도어 제어 서비스",
        "BCM_Window: 창문 제어 서비스",
        "BCM_Seat: 시트 제어 서비스",
        "BCM_SafetyBelt: 안전벨트 (V4 신규)",
        "BCM_ScreenAdjust: 스크린 조절 (V4 신규)"
    ],
    slide_num,
    code_example="""// BCM_Door API 사용 예시
class BCM_Door {
  // 도어 잠금
  Result lock(DoorPosition pos) {
    return deviceAPI.setDoorLock(pos, LOCKED);
  }
  
  // 도어 잠금 해제
  Result unlock(DoorPosition pos) {
    return deviceAPI.setDoorLock(pos, UNLOCKED);
  }
  
  // 도어 열기
  Result open(DoorPosition pos, uint8_t angle) {
    if (getSpeed() > 0) return ERROR_VEHICLE_MOVING;
    return deviceAPI.openDoor(pos, angle);
  }
  
  // 도어 상태 확인
  DoorStatus getStatus(DoorPosition pos) {
    return deviceAPI.getDoorStatus(pos);
  }
}

// 사용 예
BCM_Door door;
door.unlock(DRIVER_DOOR);
door.open(DRIVER_DOOR, 45);"""
)
slide_num += 1

# 8. TMS 도메인 API 상세
add_content_slide(
    "TMS 도메인 API 구현 예시",
    [
        "TMS_AC: 공조 시스템 제어",
        "TMS_Battery: 배터리 열관리",
        "TMS_Device: 장치 방열 제어",
        "Version 4 신규: WorkMode API 추가"
    ],
    slide_num,
    code_example="""// TMS_AC API 사용 예시
class TMS_AC {
  // 목표 온도 설정
  Result setTargetTemp(float temp) {
    if (temp < 16.0 || temp > 32.0) 
      return ERROR_OUT_OF_RANGE;
    return deviceAPI.setACTemp(temp);
  }
  
  // 팬 속도 설정
  Result setFanSpeed(uint8_t speed) {
    return deviceAPI.setFanSpeed(speed);
  }
  
  // V4 신규: 작동 모드 설정
  Result setWorkMode(TMS_Mode mode) {
    switch(mode) {
      case MODE_ECO:
        setPowerLimit(50);
        break;
      case MODE_COMFORT:
        setPowerLimit(100);
        break;
    }
    return deviceAPI.setMode(mode);
  }
}"""
)
slide_num += 1

# 9. ADAS 도메인 API 상세
add_content_slide(
    "ADAS 도메인 API 구현 예시",
    [
        "ADAS_Perception: 시각 인지",
        "ADAS_Radar: 레이더 감지",
        "ADAS_Lidar: 라이다 감지",
        "ADAS_Fusion: 센서 융합"
    ],
    slide_num,
    code_example="""// ADAS_Perception API 사용 예시
class ADAS_Perception {
  // 물체 추적
  vector<Object> getTrackObjects() {
    vector<Object> objects;
    // 카메라에서 물체 감지
    auto camObjects = camera.detect();
    for (auto& obj : camObjects) {
      objects.push_back({
        .id = obj.id,
        .type = obj.type,  // CAR, PEDESTRIAN, BICYCLE
        .position = obj.position,
        .velocity = obj.velocity,
        .confidence = obj.confidence
      });
    }
    return objects;
  }
  
  // 차선 감지
  LaneInfo getLaneline() {
    return camera.detectLanes();
  }
  
  // 신호등 인식
  TrafficLight getTrafficLight() {
    return camera.detectTrafficLight();
  }
}"""
)
slide_num += 1

# 10. Part 2: 디바이스 추상화 API 개요
add_api_table_slide(
    "Part 2: 디바이스 추상화 API - 5대 도메인 (230+ APIs)",
    [
        ["도메인", "액추에이터", "센서", "V4 신규"],
        ["BCM", "25개", "18개", "6개 모터 서비스 추가"],
        ["TMS", "12개", "8개", "매개변수 확장"],
        ["PWT", "15개", "12개", "충전포트 센서 2개"],
        ["CHS", "8개", "15개", "-"],
        ["ADAS", "3개", "8개", "-"]
    ],
    slide_num
)
slide_num += 1

# 11. Version 4 신규 모터 서비스
add_content_slide(
    "Version 4 신규 모터 서비스 API",
    [
        "단일 피드백 모터 제어",
        "2중 피드백 모터 제어",
        "3중 피드백 모터 제어",
        "가변 저항 센서 모터",
        "다단 등급 모터",
        "NTC 센서 히터"
    ],
    slide_num,
    code_example="""// V4 신규: 피드백 모터 제어
class Actr_DoubleFbMot {
  // 모터 동작 설정 (V4 확장)
  Result setOper(uint8_t dir, 
                 uint16_t dutyRat,
                 uint16_t speed) {
    MotorCmd cmd = {
      .direction = dir,
      .dutyRatio = dutyRat,  // V4: 듀티비 추가
      .targetSpeed = speed    // V4: 속도 추가
    };
    return motor.execute(cmd);
  }
  
  // V4 신규: 전압 보고
  Result ntfVolt() {
    float voltage = motor.getVoltage();
    return notify(VOLTAGE_CHANGE, voltage);
  }
  
  // V4 신규: 환경 매개변수 설정
  Result setEnvtlVal(EnvtlVal val) {
    return motor.setEnvironment(val);
  }
}"""
)
slide_num += 1

# 12. V3 vs V4 상세 비교
add_api_table_slide(
    "Version 3 → Version 4 주요 변경사항",
    [
        ["영역", "Version 3", "Version 4", "개선사항"],
        ["API 총 개수", "450개", "520개+", "15% 증가"],
        ["BCM 서비스", "29개", "31개", "SafetyBelt, ScreenAdjust 추가"],
        ["TMS API", "기본 제어", "WorkMode 추가", "3개 모드 API 신규"],
        ["모터 서비스", "기본 모터", "6종 피드백 모터", "정밀 제어 가능"],
        ["오류 등급", "4단계", "8단계", "세밀한 진단"],
        ["충전 포트", "기본 모니터링", "온도 센서 추가", "AC/DC 온도 감지"]
    ],
    slide_num
)
slide_num += 1

# 13. 중국 주요 OEM 구현 현황
add_content_slide(
    "중국 주요 OEM SDV 구현 현황",
    [
        {
            'heading': 'BYD (比亞迪)',
            'bullets': [
                'DiLink 시스템: 520개 API 중 450개 구현',
                '월간 OTA 업데이트 제공',
                '100만+ 앱 다운로드 생태계 구축',
                '1,000억 위안 SDV 투자 계획 발표'
            ]
        },
        {
            'heading': 'NIO (蔚來)',
            'bullets': [
                '중앙집중식 E/E 아키텍처 구현',
                '자체 개발 비중 80% 이상',
                'NOMI AI 어시스턴트 통합',
                '배터리 교체 시스템과 SDV 연계'
            ]
        },
        {
            'heading': 'Xiaopeng (小鵬)',
            'bullets': [
                'XPILOT 자율주행 시스템',
                'SDV 기반 스마트 콕핏',
                '분기별 주요 기능 업데이트',
                '차량-집-사무실 연결 생태계'
            ]
        }
    ],
    slide_num
)
slide_num += 1

# 14. 화웨이의 SDV 플랫폼
add_content_slide(
    "화웨이 IDVP 플랫폼과 중국 SDV 표준",
    [
        {
            'heading': '화웨이 지능형 디지털 차량 플랫폼 (IDVP)',
            'bullets': [
                '2024년 상반기 순이익 22.3억 위안 달성',
                '기업 가치 1,150억 위안 (160억 달러)',
                'SOA 기반 고도 중앙집중식 아키텍처',
                '중국 SDV 표준과 완벽 호환'
            ]
        },
        {
            'heading': '2030 전망',
            'bullets': [
                'NEV 점유율 82% 예상',
                '차량 컴퓨팅 파워 5,000 TOPS 초과',
                'L3 자율주행 30% 달성',
                '차량 네트워크 100Gbps 달성'
            ]
        }
    ],
    slide_num
)
slide_num += 1

# 15. API 호출 시퀀스 다이어그램
add_content_slide(
    "SDV API 호출 시퀀스 예시: 스마트 파킹",
    [
        "1. 사용자: '자동 주차' 버튼 터치",
        "2. HMI_Display: 주차 모드 UI 표시",
        "3. ADAS_Perception: 주차 공간 스캔",
        "4. ADAS_Fusion: 센서 데이터 융합",
        "5. VCS_Steering: 조향각 제어",
        "6. VCS_Gear: 기어 변경 (D→R→D)",
        "7. VCS_Brake: 속도 제어",
        "8. BCM_Light: 비상등 점멸",
        "9. HMI_Audio: 완료 알림음"
    ],
    slide_num,
    code_example="""// 스마트 파킹 시퀀스
async function smartParking() {
  // 1. 주차 공간 감지
  const space = await ADAS_Perception
    .getParkingSpace();
  
  // 2. 주차 가능 여부 판단
  if (space.width < 2.5 || 
      space.length < 5.0) {
    HMI_Display.showError("공간 부족");
    return;
  }
  
  // 3. 자동 주차 실행
  BCM_Light.turnOn(HAZARD_LIGHT);
  
  while (!isParked()) {
    const cmd = calculateParkingCmd();
    VCS_Steering.setAngle(cmd.angle);
    VCS_Gear.setTarget(cmd.gear);
    VCS_Brake.setTargetAx(cmd.accel);
  }
  
  // 4. 완료
  VCS_Gear.setTarget(PARK);
  BCM_Light.turnOff(HAZARD_LIGHT);
  HMI_Audio.playSound(COMPLETE);
}"""
)
slide_num += 1

# 16. 에너지 관리 시스템 API
add_content_slide(
    "EMS 도메인: 충전 및 에너지 관리",
    [
        "충전 포트 제어",
        "배터리 상태 모니터링",
        "전력 분배 최적화",
        "V2G/V2L 지원"
    ],
    slide_num,
    code_example="""// EMS_Charging API 구현
class EMS_Charging {
  // 충전 시작
  Result start(ChargeType type) {
    // V4: 충전 포트 온도 체크
    float temp = type == AC ? 
      getACTemp() : getDCTemp();
    
    if (temp > MAX_TEMP) {
      return ERROR_OVERHEATED;
    }
    
    // 충전 포트 잠금 해제
    Actr_ChrgElecLock.unlock();
    
    // 충전 시작
    return charger.startCharging({
      .type = type,
      .maxPower = getMaxPower(),
      .targetSOC = 80
    });
  }
  
  // V4 신규: 충전 포트 온도 모니터
  Result notifyACTemp() {
    return notify(AC_TEMP, sensor.getTemp());
  }
}"""
)
slide_num += 1

# 17. 실시간 차량 제어 API
add_content_slide(
    "VCS 도메인: 실시간 차량 제어",
    [
        "기어 제어 (PRND)",
        "브레이크 시스템",
        "조향 제어",
        "차량 동역학 관리"
    ],
    slide_num,
    code_example="""// VCS 통합 제어 예시
class VehicleControl {
  // 긴급 정지
  Result emergencyStop() {
    // 1. ABS 활성화
    VCS_Brake.enableABS();
    
    // 2. 최대 제동력 적용
    VCS_Brake.setTargetAx(-9.8);
    
    // 3. 안정성 제어
    VCS_Steering.enableAssist();
    
    // 4. 경고등 켜기
    BCM_Light.turnOn(HAZARD);
    
    // 5. 경고음 발생
    HMI_Audio.playAlert(EMERGENCY);
    
    return SUCCESS;
  }
  
  // 주행 모드 변경
  Result setDriveMode(DriveMode mode) {
    switch(mode) {
      case ECO:
        EMS_PowerDist.setPowerLimit(60);
        TMS_AC.setWorkMode(MODE_ECO);
        break;
      case SPORT:
        EMS_PowerDist.setPowerLimit(100);
        VCS_Steering.setSensitivity(HIGH);
        break;
    }
  }
}"""
)
slide_num += 1

# 18. 센서 융합 API
add_content_slide(
    "ADAS 센서 융합 구현",
    [
        "카메라 + 레이더 융합",
        "라이다 포인트 클라우드",
        "초음파 센서 통합",
        "IMU/GPS 데이터 동기화"
    ],
    slide_num,
    code_example="""// ADAS_Fusion API 구현
class ADAS_Fusion {
  // 멀티센서 융합
  FusedObjects getCombinedObjects() {
    // 1. 각 센서 데이터 수집
    auto camObj = camera.getObjects();
    auto radObj = radar.getObjects();
    auto lidObj = lidar.getPointCloud();
    
    // 2. 시간 동기화
    syncTimestamp(camObj, radObj, lidObj);
    
    // 3. 좌표계 변환
    transformToVehicleFrame();
    
    // 4. 칼만 필터 적용
    KalmanFilter kf;
    for (auto& obj : allObjects) {
      obj.position = kf.predict(obj);
      obj.velocity = kf.update(obj);
    }
    
    // 5. 신뢰도 계산
    calculateConfidence();
    
    return fusedObjects;
  }
}"""
)
slide_num += 1

# 19. HMI 인터랙션 API
add_content_slide(
    "HMI 도메인: 사용자 인터랙션",
    [
        "멀티 디스플레이 제어",
        "음성 인식/합성",
        "제스처 인식",
        "AR-HUD 표시"
    ],
    slide_num,
    code_example="""// HMI 통합 제어
class HMI_Controller {
  // 음성 명령 처리
  Result processVoiceCommand(string cmd) {
    auto intent = NLP.parse(cmd);
    
    switch(intent.action) {
      case "OPEN_WINDOW":
        BCM_Window.open(intent.position);
        speakResponse("창문을 엽니다");
        break;
        
      case "SET_TEMP":
        TMS_AC.setTargetTemp(intent.value);
        updateDisplay(CLIMATE_VIEW);
        break;
        
      case "NAVIGATE":
        HMI_Navigation.setDestination(
          intent.location);
        showARRoute();
        break;
    }
  }
  
  // AR-HUD 표시
  void showARRoute() {
    arHUD.project({
      .type = NAVIGATION_ARROW,
      .position = getNextTurn(),
      .distance = getDistance()
    });
  }
}"""
)
slide_num += 1

# 20. 배터리 관리 시스템
add_content_slide(
    "고전압 배터리 관리 API",
    [
        "셀 밸런싱",
        "열 관리",
        "SOC/SOH 계산",
        "고장 진단 (8단계)"
    ],
    slide_num,
    code_example="""// EMS_HVBatt 구현
class EMS_HVBatt {
  // 배터리 상태 모니터링
  BatteryStatus getStatus() {
    BatteryStatus status;
    
    // SOC 계산
    status.soc = calculateSOC();
    
    // SOH 계산
    status.soh = calculateSOH();
    
    // 셀 전압 확인
    for (int i = 0; i < CELL_COUNT; i++) {
      status.cellVoltage[i] = 
        getCellVoltage(i);
    }
    
    // V4: 8단계 오류 진단
    status.faultLevel = diagnose();
    
    // 온도 확인
    status.temp = getBatteryTemp();
    
    return status;
  }
  
  // V4: 확장된 오류 진단
  FaultLevel diagnose() {
    if (overVoltage()) return FAULT_LEVEL8;
    if (underVoltage()) return FAULT_LEVEL7;
    if (overTemp()) return FAULT_LEVEL6;
    // ... 더 세밀한 진단
    return FAULT_LEVEL0;
  }
}"""
)
slide_num += 1

# 21. 중국 SDV 시장 전망
add_content_slide(
    "중국 SDV 시장 전망 (2025-2030)",
    [
        {
            'heading': '시장 규모',
            'bullets': [
                '2024년: 2,135억 달러',
                '2030년: 1조 2,370억 달러 (CAGR 34%)',
                '중국이 아시아태평양 시장 주도',
                '글로벌 SDV 시장의 40% 차지 예상'
            ]
        },
        {
            'heading': '기술 발전',
            'bullets': [
                '차량 컴퓨팅 파워: 5,000+ TOPS',
                '네트워크 속도: 100+ Gbps',
                'L3 자율주행: 30% 보급',
                'NEV 비중: 82% 달성'
            ]
        }
    ],
    slide_num
)
slide_num += 1

# 22. 표준 채택 기업 현황
add_api_table_slide(
    "중국 SDV 표준 채택 기업 현황",
    [
        ["기업", "구현 수준", "특징", "2025 계획"],
        ["BYD", "450/520 API", "DiLink 시스템", "e4.0 플랫폼"],
        ["NIO", "80% 자체개발", "중앙집중 E/E", "NT3.0 아키텍처"],
        ["Xiaopeng", "완전 구현", "XPILOT 4.0", "도심 자율주행"],
        ["Li Auto", "부분 구현", "Li OS", "전체 구현"],
        ["Geely", "70% 구현", "GEEA 2.0", "글로벌 확대"],
        ["GWM", "개발 중", "Coffee OS", "2025 출시"]
    ],
    slide_num
)
slide_num += 1

# 23. 국제 협력 현황
add_content_slide(
    "중국 SDV 표준의 국제 협력",
    [
        {
            'heading': '글로벌 Tier-1 참여',
            'bullets': [
                'Bosch: AUTOSAR와 중국 표준 매핑',
                'Continental: 듀얼 스택 지원',
                'Aptiv: 중국향 SDV 플랫폼 개발',
                'ZF: 중국 OEM과 공동 개발'
            ]
        },
        {
            'heading': '기술 기업 협력',
            'bullets': [
                'Baidu Apollo: SDV 표준 통합',
                'Alibaba: AliOS Auto 연계',
                'Tencent: TAI 3.0 플랫폼 호환',
                'Huawei: IDVP 완전 지원'
            ]
        }
    ],
    slide_num
)
slide_num += 1

# 24. 한국의 대응 방향
add_content_slide(
    "한국의 SDV 표준 대응 전략",
    [
        {
            'heading': '현황 분석',
            'bullets': [
                '현대차그룹 ccOS 개발 중',
                'AUTOSAR 기반 접근',
                '중국 표준과의 호환성 검토 필요',
                'K-SDV 표준화 논의 시작'
            ]
        },
        {
            'heading': '제안 방향',
            'bullets': [
                '중국 표준 상세 분석 및 벤치마킹',
                'Core API + Extension 구조 채택',
                '국제 호환성 확보 (AUTOSAR/중국)',
                'MaaS/C-ITS 연계 차별화',
                '2025년 K-SDV 1.0 발표 목표'
            ]
        }
    ],
    slide_num
)
slide_num += 1

# 25. 결론
add_content_slide(
    "결론: 중국 SDV 표준의 의미",
    [
        {
            'heading': '핵심 시사점',
            'bullets': [
                '세계 최대 규모 API 표준 (520개+)',
                '빠른 반복 개발 (6개월 주기)',
                '강력한 산업계 지원과 실행력',
                '2030년 글로벌 SDV 시장 주도 가능성'
            ]
        },
        {
            'heading': '대응 필요성',
            'bullets': [
                '중국 시장 진출 시 필수 고려사항',
                '글로벌 표준 경쟁에서 주도권 확보',
                '한국형 SDV 표준 개발 가속화',
                '산학연 협력 체계 구축 시급'
            ]
        }
    ],
    slide_num
)

# 프레젠테이션 저장
output_file = 'China_SDV_Standard_Detailed_Analysis_2025.pptx'
prs.save(output_file)
print(f"중국 SDV 표준 상세 분석 프레젠테이션 생성 완료: {output_file}")
print(f"총 {slide_num} 페이지")
print("\n특징:")
print("- 2024-2025 최신 동향 반영")
print("- 520개+ API 상세 분석")
print("- 실제 코드 예시 포함")
print("- V3→V4 변경사항 상세")
print("- 중국 기업 구현 현황")
print("- 시장 전망 및 대응 전략")
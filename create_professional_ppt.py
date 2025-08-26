#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL_TYPE

def create_professional_presentation():
    """전문적인 16:9 PPT 생성"""
    # 16:9 비율 설정
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 슬라이드 1: 표지
    slide_layout = prs.slide_layouts[5]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경 색상 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 40)
    
    # 메인 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "중국 지능형 커넥티드카\n서비스 인터페이스 표준"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    
    p = subtitle_frame.add_paragraph()
    p.text = "SDV/T 001-2022 Version 4 Beta 1"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(100, 180, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 서브 정보
    info_box = slide.shapes.add_textbox(Inches(3), Inches(6), Inches(10), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.clear()
    
    p = info_frame.add_paragraph()
    p.text = "Part 1: Atomic Service API Interface\nPart 2: Device Abstraction API Interface"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.5
    
    # 하단 정보
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    
    p = footer_frame.add_paragraph()
    p.text = "한국전자기술연구원 모빌리티플랫폼연구센터 | 2025.08"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 2: 목차
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "목 차"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 40)
    p.alignment = PP_ALIGN.CENTER
    
    # 목차 내용
    contents = [
        "1. SDV 개요 및 중요성",
        "2. 중국 SDV 표준화 동향",
        "3. 표준 개발 핵심 참여사",
        "4. 4-Layer 아키텍처",
        "5. Part 1: Atomic Service API",
        "6. Part 2: Device Abstraction API",
        "7. API 연관 관계 분석",
        "8. 글로벌 경쟁 환경",
        "9. 한국형 SDV 표준 제안",
        "10. 시사점 및 결론"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(3), Inches(2), Inches(10), Inches(6))
    content_frame = content_box.text_frame
    content_frame.clear()
    
    for content in contents:
        p = content_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.line_spacing = 1.8
        p.space_after = Pt(12)
    
    # 슬라이드 3: SDV 개요
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 250)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "SDV(Software Defined Vehicle) 개요"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 50, 100)
    
    # 정의 박스
    def_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(7), Inches(2.5)
    )
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    def_box.line.color.rgb = RGBColor(100, 150, 200)
    def_box.line.width = Pt(2)
    
    text_frame = def_box.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    
    p = text_frame.add_paragraph()
    p.text = "SDV 정의"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 50, 100)
    
    p = text_frame.add_paragraph()
    p.text = "• 소프트웨어가 차량의 핵심 가치 결정\n• OTA 업데이트로 지속적 기능 개선\n• 하드웨어-소프트웨어 분리 아키텍처"
    p.font.size = Pt(18)
    p.line_spacing = 1.5
    
    # 특징 박스
    feat_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(1.5),
        Inches(7.5), Inches(2.5)
    )
    feat_box.fill.solid()
    feat_box.fill.fore_color.rgb = RGBColor(255, 240, 230)
    feat_box.line.color.rgb = RGBColor(200, 150, 100)
    feat_box.line.width = Pt(2)
    
    text_frame = feat_box.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    
    p = text_frame.add_paragraph()
    p.text = "핵심 특징"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(100, 50, 0)
    
    p = text_frame.add_paragraph()
    p.text = "• 중앙집중식 컴퓨팅 아키텍처\n• 표준화된 API 인터페이스\n• 지속적 서비스 수익 모델"
    p.font.size = Pt(18)
    p.line_spacing = 1.5
    
    # 통계 정보
    stats_y = Inches(4.5)
    for i, (num, label) in enumerate([
        ("30-34%", "연평균 성장률"),
        ("520+", "총 API 개수"),
        ("6+", "핵심 도메인")
    ]):
        stat_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(2 + i*4.5), stats_y,
            Inches(2.5), Inches(2.5)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = RGBColor(50, 150, 250)
        stat_box.line.color.rgb = RGBColor(30, 100, 200)
        stat_box.line.width = Pt(3)
        
        text_frame = stat_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 4: 중국 표준화 동향
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "중국 SDV 표준화 전략"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 0, 0)
    
    # 주도 기관
    org_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(14), Inches(1.5)
    )
    org_box.fill.solid()
    org_box.fill.fore_color.rgb = RGBColor(255, 230, 230)
    org_box.line.color.rgb = RGBColor(200, 0, 0)
    org_box.line.width = Pt(3)
    
    text_frame = org_box.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    p = text_frame.add_paragraph()
    p.text = "중국자동차공업협회(CAAM) 소프트웨어 분과 주도"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "국가 차원의 체계적 표준화 추진"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(100, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    # 전략 목표
    goals = [
        ("산업 표준화", "중국 내 산업 파편화 방지\n통일된 개발 규격 제공"),
        ("생태계 구축", "자체 기술 생태계 구축\n글로벌 주도권 확보"),
        ("비용 절감", "개발 비용 대폭 감소\n신기술 빠른 상용화")
    ]
    
    for i, (title, desc) in enumerate(goals):
        goal_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + i*5.2), Inches(3.5),
            Inches(4.8), Inches(3.5)
        )
        goal_box.fill.solid()
        goal_box.fill.fore_color.rgb = RGBColor(255, 255, 240)
        goal_box.line.color.rgb = RGBColor(200, 180, 0)
        goal_box.line.width = Pt(2)
        
        text_frame = goal_box.text_frame
        text_frame.clear()
        text_frame.margin_top = Inches(0.3)
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(150, 100, 0)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = ""
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.3
    
    # 슬라이드 5: 참여사
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 248, 252)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "표준 개발 핵심 참여사"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 80, 150)
    
    # OEM 섹션
    oem_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(7), Inches(0.6))
    oem_frame = oem_title.text_frame
    oem_frame.clear()
    
    p = oem_frame.add_paragraph()
    p.text = "주요 완성차 업체 (OEM)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 200)
    
    oems = ["BYD", "GWM", "Geely", "FAW", "SAIC", "Changan", "GAC", "Dongfeng"]
    for i, oem in enumerate(oems):
        x = Inches(0.5 + (i % 4) * 1.8)
        y = Inches(2 + (i // 4) * 1.2)
        
        oem_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(1.6), Inches(0.8)
        )
        oem_box.fill.solid()
        oem_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
        oem_box.line.color.rgb = RGBColor(100, 150, 220)
        
        text_frame = oem_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = oem
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    # Tier-1 섹션
    tier1_title = slide.shapes.add_textbox(Inches(8), Inches(1.3), Inches(7), Inches(0.6))
    tier1_frame = tier1_title.text_frame
    tier1_frame.clear()
    
    p = tier1_frame.add_paragraph()
    p.text = "주요 부품사 및 SW 기업"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 100, 0)
    
    tier1s = ["Huawei", "Bosch", "Continental", "Baidu", "Tencent", "Alibaba", "ZTE", "Xiaomi"]
    for i, tier1 in enumerate(tier1s):
        x = Inches(8 + (i % 4) * 1.8)
        y = Inches(2 + (i // 4) * 1.2)
        
        tier1_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(1.6), Inches(0.8)
        )
        tier1_box.fill.solid()
        tier1_box.fill.fore_color.rgb = RGBColor(255, 240, 230)
        tier1_box.line.color.rgb = RGBColor(220, 150, 100)
        
        text_frame = tier1_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = tier1
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    # 협력 메시지
    msg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(4.5),
        Inches(12), Inches(1.8)
    )
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = RGBColor(50, 50, 100)
    
    text_frame = msg_box.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    p = text_frame.add_paragraph()
    p.text = "산업 전체의 협력을 통한 표준 제정"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "특정 기업이 아닌 중국 자동차 산업 전체가 공동 개발"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 6: 4-Layer 아키텍처
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill  
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 250)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "SDV 서비스 소프트웨어 4계층 아키텍처"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 50, 100)
    
    # 계층 다이어그램
    layers = [
        ("애플리케이션 계층", "사용자 경험 및 차량 특화 기능 구현", RGBColor(100, 200, 100)),
        ("아토믹 서비스 계층", "Part 1: 표준화된 기능 단위 제공", RGBColor(255, 200, 100)),
        ("디바이스 추상화 계층", "Part 2: 하드웨어 제어 인터페이스", RGBColor(200, 150, 255)),
        ("기초 플랫폼 계층", "OS, 컴퓨팅 하드웨어 등 기본 환경", RGBColor(150, 200, 255))
    ]
    
    for i, (name, desc, color) in enumerate(layers):
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(1.5 + i*1.5),
            Inches(12), Inches(1.2)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        layer_box.line.color.rgb = RGBColor(80, 80, 80)
        layer_box.line.width = Pt(2)
        
        text_frame = layer_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = f"{name}: {desc}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255) if i == 3 else RGBColor(40, 40, 40)
        p.alignment = PP_ALIGN.CENTER
        
        # 화살표 추가 (마지막 레이어 제외)
        if i < 3:
            arrow = slide.shapes.add_connector(
                1, Inches(8), Inches(2.7 + i*1.5),
                Inches(8), Inches(2.9 + i*1.5)
            )
            arrow.line.color.rgb = RGBColor(100, 100, 200)
            arrow.line.width = Pt(3)
    
    # 핵심 가치
    value_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(14), Inches(1.5))
    value_frame = value_box.text_frame
    value_frame.clear()
    
    p = value_frame.add_paragraph()
    p.text = "핵심 가치: 명확한 계층 분리 | 독립적 개발 | 표준화된 인터페이스 | 하드웨어 독립성"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 80, 160)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 7: Part 1 - Atomic Service API
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 255)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "Part 1: Atomic Service API Interface"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0)
    
    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(15), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    
    p = subtitle_frame.add_paragraph()
    p.text = "아토믹 서비스 API - '기능'의 표준화 (290+ APIs)"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # 6대 도메인
    domains = [
        ("BCM", "Body Control\nModule\n차체 제어", RGBColor(100, 150, 200)),
        ("TMS", "Thermal\nManagement\n열 관리", RGBColor(200, 100, 100)),
        ("VCS", "Vehicle\nControl\n운동 제어", RGBColor(150, 200, 100)),
        ("EMS", "Energy\nManagement\n에너지 관리", RGBColor(200, 200, 100)),
        ("ADAS", "Advanced\nDriver-Assistance\n첨단 운전자 보조", RGBColor(200, 150, 200)),
        ("HMI", "Human Machine\nInterface\n사용자 인터페이스", RGBColor(150, 200, 200))
    ]
    
    for i, (code, name, color) in enumerate(domains):
        x = Inches(1 + (i % 3) * 5)
        y = Inches(2.5 + (i // 3) * 2.8)
        
        domain_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(4.5), Inches(2.3)
        )
        domain_box.fill.solid()
        domain_box.fill.fore_color.rgb = color
        domain_box.line.width = Pt(0)
        
        text_frame = domain_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = name
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2
    
    # 슬라이드 8: BCM 도메인 상세
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 252, 255)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "BCM 도메인: 차체 제어 서비스"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 80, 160)
    
    # 주요 서비스
    services = [
        ("BCM_Door", "unlock(), lock(), open(), close(), adjustPosition()"),
        ("BCM_Window", "lock(), unlock(), open(), close(), adjustPosition()"),
        ("BCM_Seat", "adjustMainXDir(), adjustBackRestAngle()"),
        ("BCM_Light", "turnOn(), turnOff() - 브레이크등, 방향지시등, 전조등"),
        ("BCM_WiperWash", "startWiping(), stopWiping(), startSprayWashing()")
    ]
    
    y_pos = Inches(2)
    for service, apis in services:
        # 서비스 명
        service_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), y_pos,
            Inches(3), Inches(0.8)
        )
        service_box.fill.solid()
        service_box.fill.fore_color.rgb = RGBColor(100, 150, 200)
        
        text_frame = service_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = service
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # API 설명
        api_box = slide.shapes.add_textbox(Inches(4.5), y_pos, Inches(10.5), Inches(0.8))
        api_frame = api_box.text_frame
        api_frame.clear()
        api_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = api_frame.add_paragraph()
        p.text = apis
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(60, 60, 60)
        
        y_pos += Inches(1.2)
    
    # 슬라이드 9: 시나리오 예시
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 248, 255)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "시나리오: '퇴근 모드' 구현"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(120, 0, 120)
    
    # 시나리오 단계
    steps = [
        ("1. 시트 및 미러 조정", 
         "BCM_Seat.adjustMainXDir(position)\nBCM_RearView.adjustXAngle(angle)"),
        ("2. 실내 환경 설정",
         "TMS_AC.setTargetTemp(22.0)\nTMS_Purifier.turnOn()"),
        ("3. 편의 기능 활성화",
         "BCM_Massage.startMassage(KNEAD, GENTLE)")
    ]
    
    for i, (step_title, code) in enumerate(steps):
        # 단계 박스
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(2 + i*2),
            Inches(14), Inches(1.5)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = RGBColor(240, 230, 255)
        step_box.line.color.rgb = RGBColor(150, 100, 200)
        step_box.line.width = Pt(2)
        
        text_frame = step_box.text_frame
        text_frame.clear()
        
        p = text_frame.add_paragraph()
        p.text = step_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 0, 100)
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.line_spacing = 1.2
    
    # 결과 메시지
    result_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(7.5),
        Inches(12), Inches(0.8)
    )
    result_box.fill.solid()
    result_box.fill.fore_color.rgb = RGBColor(100, 200, 100)
    
    text_frame = result_box.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    p = text_frame.add_paragraph()
    p.text = "여러 도메인 API를 조합하여 복잡한 시나리오 구현"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 10: Part 2 - Device Abstraction API
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 250, 245)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "Part 2: Device Abstraction API Interface"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 100, 0)
    
    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(15), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    
    p = subtitle_frame.add_paragraph()
    p.text = "디바이스 추상화 API - '장치'의 표준화 (230+ APIs)"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # 5대 도메인
    device_domains = [
        ("BCM", "Body Control\n차체 장치", RGBColor(100, 150, 200)),
        ("TMS", "Thermal Management\n열관리 장치", RGBColor(200, 100, 100)),
        ("PWT", "Powertrain\n파워트레인", RGBColor(150, 200, 100)),
        ("CHS", "Chassis\n섀시", RGBColor(200, 200, 100)),
        ("ADAS", "Advanced Driver-Assistance\nADAS 센서", RGBColor(200, 150, 200))
    ]
    
    for i, (code, name, color) in enumerate(device_domains):
        x = Inches(1 + (i % 3) * 5)
        y = Inches(2.5 + (i // 3) * 3)
        
        domain_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(4.5), Inches(2.3)
        )
        domain_box.fill.solid()
        domain_box.fill.fore_color.rgb = color
        domain_box.line.width = Pt(0)
        
        text_frame = domain_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = name
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2
    
    # 슬라이드 11: API 연관 관계
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 250, 255)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "API 계층 간 호출 흐름: '창문 열기'"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 50, 150)
    
    # 흐름도
    flow_steps = [
        ("사용자", "'창문 열기' 버튼 터치", RGBColor(100, 200, 100)),
        ("애플리케이션", "UI 이벤트 처리", RGBColor(150, 200, 150)),
        ("아토믹 서비스", "BCM_Window.open() 호출", RGBColor(255, 200, 100)),
        ("디바이스 추상화", "Actr_DoubleHallMot.setOper()", RGBColor(200, 150, 255)),
        ("하드웨어", "윈도우 모터 물리적 동작", RGBColor(150, 200, 255))
    ]
    
    for i, (layer, action, color) in enumerate(flow_steps):
        # 박스
        flow_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(1.5 + i*1.3),
            Inches(12), Inches(1)
        )
        flow_box.fill.solid()
        flow_box.fill.fore_color.rgb = color
        flow_box.line.color.rgb = RGBColor(80, 80, 80)
        flow_box.line.width = Pt(2)
        
        text_frame = flow_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = f"{layer}: {action}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 화살표 (마지막 단계 제외)
        if i < 4:
            arrow = slide.shapes.add_connector(
                1, Inches(8), Inches(2.5 + i*1.3),
                Inches(8), Inches(2.7 + i*1.3)
            )
            arrow.line.color.rgb = RGBColor(100, 100, 200)
            arrow.line.width = Pt(3)
    
    # 슬라이드 12: 한국형 표준 제안
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 252, 255)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "한국형 SDV 표준 제안 방향"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 50, 200)
    
    # 제안 사항
    proposals = [
        ("Core API + Extension",
         "핵심 API만 표준화\n확장은 Profile 형식",
         RGBColor(100, 150, 250)),
        ("국제 표준 호환성",
         "AUTOSAR, ISO 20078\nW3C VISS 호환",
         RGBColor(100, 200, 150)),
        ("보안/OTA 내장",
         "인증/권한/암호화 필수\nOTA 업데이트 API",
         RGBColor(250, 150, 100)),
        ("차량-서비스 융합",
         "MaaS, C-ITS 연계\nV2X 데이터 공유",
         RGBColor(200, 150, 250))
    ]
    
    for i, (title, desc, color) in enumerate(proposals):
        x = Inches(0.5 + (i % 2) * 8)
        y = Inches(2 + (i // 2) * 3)
        
        prop_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(7), Inches(2.5)
        )
        prop_box.fill.solid()
        prop_box.fill.fore_color.rgb = color
        
        text_frame = prop_box.text_frame
        text_frame.clear()
        text_frame.margin_top = Inches(0.3)
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = ""
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.3
    
    # 슬라이드 13: 시사점 및 결론
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 30, 50)
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    
    p = title_frame.add_paragraph()
    p.text = "결론 및 시사점"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 핵심 메시지
    core_msg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2),
        Inches(14), Inches(2)
    )
    core_msg.fill.solid()
    core_msg.fill.fore_color.rgb = RGBColor(50, 100, 200)
    
    text_frame = core_msg.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    p = text_frame.add_paragraph()
    p.text = "중국 SDV 표준은 기술 혁신과 생태계 주도권을 동시에 추구"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "글로벌 자동차 산업의 새로운 경쟁 축 형성"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 시사점 리스트
    implications = [
        "✓ 중국 시장 진출 시 표준 준수 필수화 가능성",
        "✓ 국내 K-SDV 표준화 가속화 필요",
        "✓ 완성차-IT-부품업체 협력 생태계 구축",
        "✓ 고부가가치 소프트웨어 시장 대비"
    ]
    
    for i, imp in enumerate(implications):
        imp_box = slide.shapes.add_textbox(Inches(2), Inches(4.5 + i*0.8), Inches(12), Inches(0.7))
        imp_frame = imp_box.text_frame
        imp_frame.clear()
        
        p = imp_frame.add_paragraph()
        p.text = imp
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 마지막 슬라이드: Q&A
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 40)
    
    # Q&A
    qa_box = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(10), Inches(2))
    qa_frame = qa_box.text_frame
    qa_frame.clear()
    
    p = qa_frame.add_paragraph()
    p.text = "Q & A"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 감사 메시지
    thanks_box = slide.shapes.add_textbox(Inches(3), Inches(6), Inches(10), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.clear()
    
    p = thanks_frame.add_paragraph()
    p.text = "감사합니다"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(150, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 하단 정보
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    
    p = footer_frame.add_paragraph()
    p.text = "한국전자기술연구원 모빌리티플랫폼연구센터"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    return prs

# 프레젠테이션 생성 및 저장
prs = create_professional_presentation()
prs.save('/home/kim/github-sdv/China_SDV_Standard_Analysis_v4_Professional.pptx')
print("전문적인 16:9 프레젠테이션이 생성되었습니다: China_SDV_Standard_Analysis_v4_Professional.pptx")
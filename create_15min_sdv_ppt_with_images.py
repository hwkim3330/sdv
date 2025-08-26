#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
중국 SDV 표준 15분 발표용 PPT (이미지 포함)
간결하고 임팩트 있는 구성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 프레젠테이션 생성 (16:9 비율)
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 색상 정의
CHINA_RED = RGBColor(238, 28, 37)
DARK_BLUE = RGBColor(0, 32, 96)
LIGHT_BLUE = RGBColor(218, 238, 243)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)

def add_title_slide_with_image():
    """타이틀 슬라이드 with 이미지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경 이미지 (SDV 관련)
    try:
        img_path = "img/230517_dSPACE_SDV_Master_v01.png"
        if os.path.exists(img_path):
            pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), 
                                          width=Inches(16), height=Inches(9))
            # 이미지를 맨 뒤로 보내기
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            # 투명도 설정
            pic.transparency = 0.5
    except:
        pass
    
    # 반투명 오버레이
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2),
        Inches(16), Inches(5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = WHITE
    overlay.fill.transparency = 0.3
    overlay.line.fill.background()
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
    title = title_box.text_frame
    title.text = "중국 SDV 표준 분석"
    title.paragraphs[0].font.size = Pt(60)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = CHINA_RED
    title.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 서브타이틀
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
    subtitle = subtitle_box.text_frame
    subtitle.text = "SDV/T 001-2022 Version 4 Beta 1"
    subtitle.paragraphs[0].font.size = Pt(32)
    subtitle.paragraphs[0].font.color.rgb = DARK_BLUE
    subtitle.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # KETI 로고 위치
    subtitle2_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(14), Inches(1))
    subtitle2 = subtitle2_box.text_frame
    subtitle2.text = "한국전자기술연구원 (KETI)"
    subtitle2.paragraphs[0].font.size = Pt(20)
    subtitle2.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_slide_with_image(title_text, content, slide_num, img_file=None, img_position="right"):
    """이미지가 포함된 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(14), Inches(0.6))
    title = title_box.text_frame
    title.text = title_text
    title.paragraphs[0].font.size = Pt(28)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    
    # 이미지 추가
    if img_file and os.path.exists(f"img/{img_file}"):
        try:
            if img_position == "right":
                # 오른쪽에 이미지
                pic = slide.shapes.add_picture(f"img/{img_file}", 
                                              Inches(9), Inches(1.5), 
                                              width=Inches(6.5))
                content_width = Inches(8)
            elif img_position == "bottom":
                # 하단에 이미지
                pic = slide.shapes.add_picture(f"img/{img_file}", 
                                              Inches(2), Inches(5), 
                                              width=Inches(12))
                content_width = Inches(14)
            elif img_position == "center":
                # 중앙에 크게
                pic = slide.shapes.add_picture(f"img/{img_file}", 
                                              Inches(2), Inches(1.5), 
                                              width=Inches(12))
                content_width = Inches(14)
        except:
            content_width = Inches(14)
    else:
        content_width = Inches(14)
    
    # 컨텐츠
    if img_position != "center":
        content_top = Inches(1.5)
        for item in content:
            if isinstance(item, dict):
                # 헤딩
                if 'heading' in item:
                    heading_box = slide.shapes.add_textbox(Inches(0.5), content_top, content_width, Inches(0.6))
                    heading = heading_box.text_frame
                    heading.text = item['heading']
                    heading.paragraphs[0].font.size = Pt(22)
                    heading.paragraphs[0].font.bold = True
                    heading.paragraphs[0].font.color.rgb = CHINA_RED
                    content_top += Inches(0.6)
                
                # 불릿
                if 'bullets' in item:
                    for bullet in item['bullets']:
                        bullet_box = slide.shapes.add_textbox(Inches(1), content_top, content_width-Inches(0.5), Inches(0.5))
                        bullet_text = bullet_box.text_frame
                        bullet_text.text = f"• {bullet}"
                        bullet_text.paragraphs[0].font.size = Pt(18)
                        bullet_text.paragraphs[0].font.color.rgb = BLACK
                        content_top += Inches(0.5)
            else:
                # 일반 텍스트
                text_box = slide.shapes.add_textbox(Inches(0.5), content_top, content_width, Inches(0.5))
                text = text_box.text_frame
                text.text = f"• {item}"
                text.paragraphs[0].font.size = Pt(20)
                text.paragraphs[0].font.color.rgb = BLACK
                content_top += Inches(0.6)
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    page = page_box.text_frame
    page.text = str(slide_num)
    page.paragraphs[0].font.size = Pt(12)
    page.paragraphs[0].font.color.rgb = GRAY
    page.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_company_logos_slide(slide_num):
    """기업 로고 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(14), Inches(0.6))
    title = title_box.text_frame
    title.text = "중국 SDV 표준 참여 주요 기업"
    title.paragraphs[0].font.size = Pt(28)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    
    # 설명 텍스트
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    desc = desc_box.text_frame
    desc.text = "60개 이상 기업이 참여하여 520개+ API 표준 개발"
    desc.paragraphs[0].font.size = Pt(20)
    desc.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 로고들 추가
    logos = [
        ("BYD_Brazil_company.png", Inches(1), Inches(2.5), "BYD"),
        ("geely-2023-logo-png_seeklogo-466890.png", Inches(4.5), Inches(2.5), "Geely"),
        ("Huawei_Standard_logo.svg.png", Inches(8), Inches(2.5), "Huawei"),
    ]
    
    for logo_file, x, y, name in logos:
        try:
            if os.path.exists(f"img/{logo_file}"):
                pic = slide.shapes.add_picture(f"img/{logo_file}", x, y, width=Inches(3))
                # 회사명 추가
                name_box = slide.shapes.add_textbox(x, y + Inches(1.8), Inches(3), Inches(0.5))
                name_text = name_box.text_frame
                name_text.text = name
                name_text.paragraphs[0].font.size = Pt(16)
                name_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        except:
            pass
    
    # 통계 박스들
    stats = [
        ("520+", "총 API 수"),
        ("6+5", "도메인 수"),
        ("60+", "참여 기업"),
        ("2030", "목표 연도")
    ]
    
    stat_left = Inches(1)
    for stat_num, stat_desc in stats:
        stat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            stat_left, Inches(5.5),
            Inches(3.5), Inches(1.5)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = LIGHT_BLUE
        
        stat_text = stat_box.text_frame
        stat_text.text = f"{stat_num}\n{stat_desc}"
        for para in stat_text.paragraphs:
            para.alignment = PP_ALIGN.CENTER
        stat_text.paragraphs[0].font.size = Pt(28)
        stat_text.paragraphs[0].font.bold = True
        stat_text.paragraphs[0].font.color.rgb = CHINA_RED
        stat_text.paragraphs[1].font.size = Pt(16)
        
        stat_left += Inches(3.8)
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    page = page_box.text_frame
    page.text = str(slide_num)
    page.paragraphs[0].font.size = Pt(12)
    page.paragraphs[0].font.color.rgb = GRAY

def add_architecture_slide_with_image(slide_num):
    """아키텍처 다이어그램 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(14), Inches(0.6))
    title = title_box.text_frame
    title.text = "중국 SDV 4계층 아키텍처"
    title.paragraphs[0].font.size = Pt(28)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    
    # 아키텍처 이미지 추가 시도
    arch_images = [
        "F184_[그림1] 오비고 SDV 플랫폼 아키텍처.png",
        "F185_[그림2] SDV용 자율주행 SW 구조도.png",
        "hyundai-ces-sdv-ee-architecture-pc.avif",
        "virtualization-microservice-architecture-for-software-defined-vehicles-an-evaluation-and-exploration-0.png"
    ]
    
    image_added = False
    for img in arch_images:
        if os.path.exists(f"img/{img}"):
            try:
                pic = slide.shapes.add_picture(f"img/{img}", 
                                              Inches(8), Inches(1.5), 
                                              width=Inches(7))
                image_added = True
                break
            except:
                pass
    
    # 4계층 설명
    layers = [
        ("Layer 1: 애플리케이션", "OEM Apps, 3rd Party Apps", LIGHT_BLUE),
        ("Layer 2: 아토믹 서비스 (Part 1)", "BCM, TMS, VCS, EMS, ADAS, HMI", RGBColor(230, 255, 230)),
        ("Layer 3: 디바이스 추상화 (Part 2)", "Actuators, Sensors, ECUs", RGBColor(230, 230, 255)),
        ("Layer 4: 기초 플랫폼", "Linux, QNX, Android Auto", RGBColor(240, 240, 240))
    ]
    
    layer_top = Inches(1.5)
    for layer_name, layer_desc, color in layers:
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), layer_top,
            Inches(7), Inches(1.3)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        
        layer_text = layer_box.text_frame
        layer_text.text = f"{layer_name}\n{layer_desc}"
        layer_text.paragraphs[0].font.size = Pt(18)
        layer_text.paragraphs[0].font.bold = True
        layer_text.paragraphs[1].font.size = Pt(14)
        for para in layer_text.paragraphs:
            para.alignment = PP_ALIGN.CENTER
        
        layer_top += Inches(1.5)
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    page = page_box.text_frame
    page.text = str(slide_num)
    page.paragraphs[0].font.size = Pt(12)
    page.paragraphs[0].font.color.rgb = GRAY

def add_api_comparison_table(slide_num):
    """API 비교 테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(14), Inches(0.6))
    title = title_box.text_frame
    title.text = "Version 3 → Version 4 핵심 변경사항"
    title.paragraphs[0].font.size = Pt(28)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = WHITE
    
    # 테이블 추가
    table_data = [
        ["구분", "Version 3", "Version 4", "개선율"],
        ["총 API 수", "450개", "520개+", "+15%"],
        ["BCM 서비스", "29개", "31개", "+7%"],
        ["신규 모터", "기본 모터", "6종 피드백", "+600%"],
        ["오류 진단", "4단계", "8단계", "+100%"],
        ["온도 센서", "기본", "AC/DC 분리", "+100%"]
    ]
    
    rows = len(table_data)
    cols = len(table_data[0])
    
    table = slide.shapes.add_table(rows, cols, 
                                   Inches(1), Inches(1.5), 
                                   Inches(14), Inches(4)).table
    
    # 테이블 스타일링
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(table_data[i][j])
            
            if i == 0:  # 헤더
                cell.fill.solid()
                cell.fill.fore_color.rgb = CHINA_RED
                para = cell.text_frame.paragraphs[0]
                para.font.color.rgb = WHITE
                para.font.bold = True
                para.font.size = Pt(18)
            else:
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(16)
                if j == 3:  # 개선율 컬럼
                    para.font.color.rgb = CHINA_RED
                    para.font.bold = True
                else:
                    para.font.color.rgb = BLACK
            
            para.alignment = PP_ALIGN.CENTER
    
    # 하단 설명
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(1.5))
    desc = desc_box.text_frame
    desc.text = "✓ 6개 신규 모터 서비스: SingleFbMot, DoubleFbMot, TripleFbMot, SlideRMot, GradedMot, Heatr\n"
    desc.text += "✓ BCM 신규: SafetyBelt (안전벨트), ScreenAdjust (스크린 조절)\n"
    desc.text += "✓ TMS 개선: WorkMode API 3개 추가, 매개변수 확장"
    for para in desc.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = DARK_BLUE
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.5))
    page = page_box.text_frame
    page.text = str(slide_num)
    page.paragraphs[0].font.size = Pt(12)
    page.paragraphs[0].font.color.rgb = GRAY

# 슬라이드 생성
slide_num = 1

# 1. 타이틀 슬라이드
add_title_slide_with_image()
slide_num += 1

# 2. 중국 SDV 시장 현황
add_slide_with_image(
    "중국 SDV 시장 현황 2025",
    [
        {
            'heading': '시장 규모',
            'bullets': [
                '2024년: 2,135억 달러',
                '2030년: 1조 2,370억 달러 (CAGR 34%)',
                '글로벌 시장 40% 점유 예상'
            ]
        },
        {
            'heading': '2025년 목표',
            'bullets': [
                'NEV 판매: 1,600만대 (YoY +24.4%)',
                'L3 자율주행: 30% 보급',
                '차량 컴퓨팅: 5,000 TOPS'
            ]
        }
    ],
    slide_num,
    "sqXZYbvJ2rUg8RGPUDZleXzTT78.jpg",
    "right"
)
slide_num += 1

# 3. SDV 표준 개요
add_slide_with_image(
    "중국 SDV 표준 개요",
    [
        {
            'heading': 'SDV/T 001-2022 Version 4 Beta 1',
            'bullets': [
                '발표: 2022년 12월',
                '주관: CAAM 소프트웨어 분과',
                '참여: 60개+ 기업'
            ]
        },
        {
            'heading': '표준 구성',
            'bullets': [
                'Part 1: 아토믹 서비스 API (290+)',
                'Part 2: 디바이스 추상화 API (230+)',
                '총 520개+ API 정의'
            ]
        }
    ],
    slide_num,
    "img.jpg",
    "right"
)
slide_num += 1

# 4. 참여 기업 로고
add_company_logos_slide(slide_num)
slide_num += 1

# 5. 4계층 아키텍처
add_architecture_slide_with_image(slide_num)
slide_num += 1

# 6. Part 1 아토믹 서비스 API
add_slide_with_image(
    "Part 1: 아토믹 서비스 API (6대 도메인)",
    [
        "BCM: 차체 제어 (31개 서비스)",
        "TMS: 열 관리 (8개 서비스)",
        "VCS: 차량 제어 (12개 서비스)",
        "EMS: 에너지 관리 (15개 서비스)",
        "ADAS: 지능형 주행 (18개 서비스)",
        "HMI: 사용자 인터페이스 (10개 서비스)"
    ],
    slide_num,
    "P12(0).jpg",
    "right"
)
slide_num += 1

# 7. Part 2 디바이스 추상화 API
add_slide_with_image(
    "Part 2: 디바이스 추상화 API (5대 도메인)",
    [
        "BCM: 액추에이터 25개, 센서 18개",
        "TMS: 액추에이터 12개, 센서 8개",
        "PWT: 액추에이터 15개, 센서 12개",
        "CHS: 액추에이터 8개, 센서 15개",
        "ADAS: 액추에이터 3개, 센서 8개"
    ],
    slide_num,
    "P9(0).jpg",
    "right"
)
slide_num += 1

# 8. V3→V4 변경사항
add_api_comparison_table(slide_num)
slide_num += 1

# 9. BYD 구현 사례
add_slide_with_image(
    "중국 OEM 구현 사례: BYD",
    [
        {
            'heading': 'DiLink 시스템',
            'bullets': [
                '520개 API 중 450개 구현',
                '월간 OTA 업데이트',
                '100만+ 앱 생태계'
            ]
        },
        {
            'heading': '2025 Xuanji 아키텍처',
            'bullets': [
                'God\'s Eye ADAS (100만대+)',
                '1,000억 위안 투자',
                '21개 모델 스마트 드라이빙'
            ]
        }
    ],
    slide_num,
    "01HBWV5X4T5EF6P8SYD8JGDKEV.jpg",
    "right"
)
slide_num += 1

# 10. 화웨이 IDVP
add_slide_with_image(
    "화웨이 IDVP 플랫폼",
    [
        {
            'heading': '2024년 성과',
            'bullets': [
                '순이익: 22.3억 위안',
                '기업가치: 160억 달러',
                'SDV 표준 완벽 호환'
            ]
        },
        {
            'heading': '2030 전망',
            'bullets': [
                'NEV 점유율: 82%',
                '차량 컴퓨팅: 5,000+ TOPS',
                'L3 자율주행: 30%'
            ]
        }
    ],
    slide_num,
    "Huawei_Standard_logo.svg.png",
    "right"
)
slide_num += 1

# 11. 한국 대응 전략
add_slide_with_image(
    "한국의 SDV 대응 전략",
    [
        {
            'heading': '현황',
            'bullets': [
                '현대차 ccOS 개발 중',
                'AUTOSAR 기반 접근',
                'K-SDV 표준화 논의'
            ]
        },
        {
            'heading': '제안',
            'bullets': [
                '중국 표준 벤치마킹',
                'Core API + Extension 구조',
                'MaaS/C-ITS 차별화',
                '2025년 K-SDV 1.0 목표'
            ]
        }
    ],
    slide_num,
    "hyundai-ces-sdv-hpvc-pc.avif",
    "right"
)
slide_num += 1

# 12. 핵심 시사점
add_slide_with_image(
    "핵심 시사점",
    [
        "✓ 세계 최대 규모 API 표준 (520개+)",
        "✓ 빠른 반복 개발 (6개월 주기)",
        "✓ 강력한 산업계 지원",
        "✓ 2030년 글로벌 시장 주도 가능성",
        "",
        "⚡ 한국의 신속한 대응 필요",
        "⚡ 글로벌 호환성 확보 중요",
        "⚡ 차별화된 K-SDV 전략 수립"
    ],
    slide_num
)
slide_num += 1

# 13. Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 배경색
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_BLUE

# Q&A 텍스트
qa_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
qa = qa_box.text_frame
qa.text = "Q & A"
qa.paragraphs[0].font.size = Pt(72)
qa.paragraphs[0].font.bold = True
qa.paragraphs[0].font.color.rgb = DARK_BLUE
qa.paragraphs[0].alignment = PP_ALIGN.CENTER

# 연락처
contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(1.5))
contact = contact_box.text_frame
contact.text = "감사합니다\n\n한국전자기술연구원 (KETI)\n모빌리티플랫폼연구센터"
for para in contact.paragraphs:
    para.font.size = Pt(20)
    para.alignment = PP_ALIGN.CENTER

# 프레젠테이션 저장
output_file = 'China_SDV_Standard_15min_Presentation_with_Images.pptx'
prs.save(output_file)
print(f"✅ 15분 발표용 PPT 생성 완료: {output_file}")
print(f"📊 총 {slide_num} 페이지")
print("\n특징:")
print("- 15분 발표에 최적화 (13페이지)")
print("- 실제 이미지 포함")
print("- 기업 로고 삽입")
print("- 시각적 아키텍처 다이어그램")
print("- 핵심 메시지 중심 구성")
print("\n발표 시간 가이드:")
print("- 슬라이드 1-2: 2분 (도입)")
print("- 슬라이드 3-5: 3분 (표준 개요)")
print("- 슬라이드 6-8: 4분 (API 상세)")
print("- 슬라이드 9-10: 3분 (기업 사례)")
print("- 슬라이드 11-12: 3분 (한국 전략)")
print("- 슬라이드 13: Q&A")